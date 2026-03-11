"""
generar_documentos.py — v1.0
Genera Informe Word y Presentación PowerPoint del Proyecto WMS Egakat.
Uso: py generar_documentos.py
"""

import sys, os, subprocess
from datetime import datetime

sys.stdout.reconfigure(encoding="utf-8")

# ─── Instalar dependencias si no están ──────────────────────────────────
print("Verificando dependencias...")
for pkg in ["python-docx", "python-pptx"]:
    subprocess.run([sys.executable, "-m", "pip", "install", pkg, "-q"], check=False)

# ─── Imports docx ────────────────────────────────────────────────────────
from docx import Document
from docx.shared import Inches as WInches, Pt as WPt, RGBColor as WRGBColor, Cm
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.oxml.ns import qn
from docx.oxml import OxmlElement

# ─── Imports pptx ────────────────────────────────────────────────────────
from pptx import Presentation
from pptx.util import Inches as PInches, Pt as PPt, Emu
from pptx.dml.color import RGBColor as PRGBColor
from pptx.enum.text import PP_ALIGN

# ─── Constantes ──────────────────────────────────────────────────────────
HOY          = datetime.now().strftime("%d de %B de %Y")
CARPETA      = r"C:\ClaudeWork"

# Paleta Word (WRGBColor)
W_AZUL_OSC  = WRGBColor(0x1F, 0x4E, 0x79)
W_AZUL_MED  = WRGBColor(0x2E, 0x75, 0xB6)
W_VERDE     = WRGBColor(0x70, 0xAD, 0x47)
W_NARANJA   = WRGBColor(0xED, 0x7D, 0x31)
W_BLANCO    = WRGBColor(0xFF, 0xFF, 0xFF)
W_GRIS      = WRGBColor(0x59, 0x59, 0x59)
W_ROJO      = WRGBColor(0xC0, 0x00, 0x00)

# Paleta PPT (PRGBColor)
P_AZUL_OSC  = PRGBColor(0x1F, 0x4E, 0x79)
P_AZUL_MED  = PRGBColor(0x2E, 0x75, 0xB6)
P_VERDE     = PRGBColor(0x70, 0xAD, 0x47)
P_NARANJA   = PRGBColor(0xED, 0x7D, 0x31)
P_BLANCO    = PRGBColor(0xFF, 0xFF, 0xFF)
P_GRIS      = PRGBColor(0x80, 0x80, 0x80)
P_GRIS_CLR  = PRGBColor(0xF2, 0xF2, 0xF2)
P_NEGRO     = PRGBColor(0x33, 0x33, 0x33)


# ════════════════════════════════════════════════════════════════════════
#  HELPERS WORD
# ════════════════════════════════════════════════════════════════════════

def w_cell_bg(cell, hex6):
    tc   = cell._tc
    tcPr = tc.get_or_add_tcPr()
    shd  = OxmlElement("w:shd")
    shd.set(qn("w:val"),   "clear")
    shd.set(qn("w:color"), "auto")
    shd.set(qn("w:fill"),  hex6)
    tcPr.append(shd)


def w_heading(doc, texto, level, color=None):
    p = doc.add_heading(texto, level=level)
    if color:
        for run in p.runs:
            run.font.color.rgb = color
    return p


def w_caja(doc, texto, bg_hex="DEEAF1", text_color=None, bold=False, font_size=10):
    t    = doc.add_table(rows=1, cols=1)
    t.style = "Table Grid"
    cell = t.cell(0, 0)
    w_cell_bg(cell, bg_hex)
    p    = cell.paragraphs[0]
    run  = p.add_run(texto)
    run.font.size = WPt(font_size)
    run.font.bold = bold
    if text_color:
        run.font.color.rgb = text_color
    doc.add_paragraph()
    return t


# ════════════════════════════════════════════════════════════════════════
#  HELPERS PPT
# ════════════════════════════════════════════════════════════════════════

BLANK_LAYOUT = None   # se asigna al iniciar la presentación


def p_bg(slide, color):
    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = color


def p_rect(slide, left, top, w, h, fill, border=None):
    shape = slide.shapes.add_shape(
        1,
        PInches(left), PInches(top), PInches(w), PInches(h)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if border:
        shape.line.color.rgb = border
    else:
        shape.line.color.rgb = fill   # invisible (same as fill)
    return shape


def p_txt(slide, texto, left, top, w, h,
          size=12, bold=False, italic=False,
          color=None, align=PP_ALIGN.LEFT, wrap=True):
    tb = slide.shapes.add_textbox(
        PInches(left), PInches(top), PInches(w), PInches(h)
    )
    tf = tb.text_frame
    tf.word_wrap = wrap
    pg = tf.paragraphs[0]
    pg.alignment = align
    run = pg.add_run()
    run.text = texto
    run.font.size  = PPt(size)
    run.font.bold  = bold
    run.font.italic = italic
    if color:
        run.font.color.rgb = color
    return tb


# ════════════════════════════════════════════════════════════════════════
#  GENERAR WORD
# ════════════════════════════════════════════════════════════════════════

def generar_word():
    doc = Document()

    # Márgenes
    for sec in doc.sections:
        sec.top_margin    = Cm(2.0)
        sec.bottom_margin = Cm(2.0)
        sec.left_margin   = Cm(2.5)
        sec.right_margin  = Cm(2.5)

    # ── PORTADA ──────────────────────────────────────────────────────
    doc.add_paragraph()
    doc.add_paragraph()

    # Banda superior
    t0 = doc.add_table(rows=1, cols=1); t0.style = "Table Grid"
    c0 = t0.cell(0, 0); w_cell_bg(c0, "1F4E79")
    p0 = c0.paragraphs[0]; p0.alignment = WD_ALIGN_PARAGRAPH.CENTER
    r0 = p0.add_run("  EGAKAT SPA   ·   CONTROL DE GESTIÓN Y MEJORA CONTINUA  ")
    r0.font.color.rgb = W_BLANCO; r0.font.bold = True; r0.font.size = WPt(11)

    doc.add_paragraph()
    doc.add_paragraph()

    pt = doc.add_paragraph(); pt.alignment = WD_ALIGN_PARAGRAPH.CENTER
    rt = pt.add_run("AUTOMATIZACIÓN DE REPORTES WMS")
    rt.font.size = WPt(28); rt.font.bold = True; rt.font.color.rgb = W_AZUL_OSC

    doc.add_paragraph()

    ps = doc.add_paragraph(); ps.alignment = WD_ALIGN_PARAGRAPH.CENTER
    rs = ps.add_run("Transformación Digital del Control de Gestión")
    rs.font.size = WPt(16); rs.font.color.rgb = W_AZUL_MED

    doc.add_paragraph()
    doc.add_paragraph()

    t1 = doc.add_table(rows=1, cols=1); t1.style = "Table Grid"
    c1 = t1.cell(0, 0); w_cell_bg(c1, "DEEAF1")
    p1 = c1.paragraphs[0]; p1.alignment = WD_ALIGN_PARAGRAPH.CENTER
    r1 = p1.add_run(
        f"Autor: Sócrates Cabral\n"
        f"Head of Control Management & Continuous Improvement\n"
        f"{HOY}"
    )
    r1.font.size = WPt(12); r1.font.color.rgb = W_AZUL_OSC

    doc.add_page_break()

    # ── 1. RESUMEN EJECUTIVO ─────────────────────────────────────────
    w_heading(doc, "1. RESUMEN EJECUTIVO", 1, W_AZUL_OSC)
    doc.add_paragraph(
        "Este informe documenta el diseño, implementación y resultados del proyecto de "
        "automatización de reportes del Sistema de Gestión de Almacenes (WMS) de Egakat SPA. "
        "El proyecto transforma un proceso manual, repetitivo y propenso a errores en una "
        "solución completamente automatizada que garantiza datos precisos, consistentes y "
        "disponibles a primera hora del día laboral."
    )
    doc.add_paragraph()

    # KPIs
    kpi_tabla = doc.add_table(rows=2, cols=4)
    kpi_tabla.style = "Table Grid"
    kpi_h = ["Reportes\nAutomatizados", "Centros de\nDistribución", "Empresas\nClientes", "Horas Ahorradas\npor Día"]
    kpi_v = ["27", "3", "16", "~2 horas"]
    for i, (h_t, v_t) in enumerate(zip(kpi_h, kpi_v)):
        hc = kpi_tabla.cell(0, i); vc = kpi_tabla.cell(1, i)
        w_cell_bg(hc, "1F4E79"); w_cell_bg(vc, "DEEAF1")
        ph = hc.paragraphs[0]; ph.alignment = WD_ALIGN_PARAGRAPH.CENTER
        rh = ph.add_run(h_t); rh.font.color.rgb = W_BLANCO; rh.font.bold = True; rh.font.size = WPt(9)
        pv = vc.paragraphs[0]; pv.alignment = WD_ALIGN_PARAGRAPH.CENTER
        rv = pv.add_run(v_t); rv.font.bold = True; rv.font.size = WPt(16); rv.font.color.rgb = W_AZUL_OSC

    doc.add_paragraph()
    doc.add_paragraph(
        "La automatización cubre tres módulos críticos: Stock WMS Semanal, Staging IN/OUT y "
        "Consulta de Posiciones — los tres pilares de información que alimentan el ecosistema "
        "de reportería y toma de decisiones de Egakat SPA."
    )
    doc.add_page_break()

    # ── 2. EL DESAFÍO ────────────────────────────────────────────────
    w_heading(doc, "2. EL DESAFÍO: ANTES DE LA AUTOMATIZACIÓN", 1, W_AZUL_OSC)
    doc.add_paragraph(
        "Cada día hábil, el equipo de Control de Gestión debía ejecutar manualmente un proceso "
        "secuencial para obtener la información operacional del WMS. Este proceso consumía "
        "recursos valiosos y presentaba riesgos significativos para la calidad del dato."
    )
    doc.add_paragraph()

    antes_t = doc.add_table(rows=5, cols=2)
    antes_t.style = "Table Grid"
    antes_t.cell(0, 0).merge(antes_t.cell(0, 1))
    w_cell_bg(antes_t.cell(0, 0), "C00000")
    p_m = antes_t.cell(0, 0).paragraphs[0]; p_m.alignment = WD_ALIGN_PARAGRAPH.CENTER
    r_m = p_m.add_run("PROCESO MANUAL — SITUACIÓN ANTERIOR")
    r_m.font.color.rgb = W_BLANCO; r_m.font.bold = True

    antes_data = [
        ("Tiempo diario",             "~90-120 minutos de trabajo manual repetitivo"),
        ("Consistencia",              "Variable — dependía del operador del momento"),
        ("Disponibilidad de datos",   "Tardía — disponible a media mañana (9:30-10:00 AM)"),
        ("Riesgo de error",           "Alto — nombres de archivo inconsistentes, omisiones frecuentes"),
    ]
    for i, (c1, c2) in enumerate(antes_data, 1):
        w_cell_bg(antes_t.cell(i, 0), "FCE4D6")
        r1 = antes_t.cell(i, 0).paragraphs[0].add_run(c1)
        r1.font.bold = True; r1.font.size = WPt(10)
        antes_t.cell(i, 1).paragraphs[0].add_run(c2).font.size = WPt(10)

    doc.add_paragraph()
    w_caja(doc,
        "El proceso manual no solo consumía tiempo: cada descarga representaba un punto de "
        "falla potencial. Un archivo en la carpeta equivocada o con nombre distinto rompía "
        "las conexiones de Power BI, generando reportes desactualizados para la operación.",
        bg_hex="FFF2CC")

    doc.add_paragraph()
    w_heading(doc, "Principales Puntos de Dolor", 2, W_AZUL_MED)
    pain_pts = [
        "~2 horas de trabajo calificado consumidas en tareas repetitivas sin valor analítico",
        "Con 27 reportes diarios, la probabilidad de omitir uno era alta",
        "Inconsistencia en nombres de archivos → ruptura de conexiones Power BI / Power Query",
        "Datos disponibles a media mañana → decisiones operativas postergadas",
        "Dependencia de una persona específica → riesgo ante vacaciones o ausencias",
        "Sin registro auditado de cuándo se descargó cada reporte ni de errores ocurridos",
    ]
    for pp in pain_pts:
        p = doc.add_paragraph(style="List Bullet")
        p.add_run(pp).font.size = WPt(10)

    doc.add_page_break()

    # ── 3. LA SOLUCIÓN ────────────────────────────────────────────────
    w_heading(doc, "3. LA SOLUCIÓN: AUTOMATIZACIÓN INTEGRAL", 1, W_AZUL_OSC)
    doc.add_paragraph(
        "El proyecto implementa tres módulos de automatización que trabajan de forma coordinada, "
        "ejecutándose en secuencia cada día hábil a las 8:00 AM vía el Programador de Tareas de "
        "Windows. Cada módulo es robusto ante errores, idempotente y deja trazabilidad completa."
    )
    doc.add_paragraph()

    modulos_w = [
        ("MÓDULO 1", "Stock WMS Semanal",        "2E75B6",
         "Descarga el inventario completo de los 3 centros de distribución activos "
         "(Quilicura, Pudahuel, Pudahuel Unitario). Genera 3 archivos Excel con nombres "
         "estandarizados para consumo directo en Power BI.",
         "3 reportes  |  3 CDs  |  ~5 minutos"),
        ("MÓDULO 2", "Staging IN/OUT",            "ED7D31",
         "Descarga el reporte 'Consulta Stock con Staging In y Out' para 16 empresas "
         "clientes distribuidas en 3 sesiones WMS independientes. Los archivos CSV se "
         "guardan con el nombre original del sistema para trazabilidad completa.",
         "16 reportes  |  16 clientes  |  ~10 minutos"),
        ("MÓDULO 3", "Consulta de Posiciones",    "70AD47",
         "Descarga los 8 reportes de posiciones (ocupadas y libres) de los 4 depósitos. "
         "Utiliza nombres de archivo fijos que Power Query consume directamente sin "
         "ninguna configuración adicional.",
         "8 reportes  |  4 depósitos  |  ~4 minutos"),
    ]

    for cod, nombre, color_hex, desc, stats in modulos_w:
        tm = doc.add_table(rows=1, cols=2)
        tm.style = "Table Grid"
        cl = tm.cell(0, 0); cr = tm.cell(0, 1)
        w_cell_bg(cl, color_hex); w_cell_bg(cr, "F2F2F2")
        pl = cl.paragraphs[0]; pl.alignment = WD_ALIGN_PARAGRAPH.CENTER
        rl1 = pl.add_run(f"{cod}\n"); rl1.font.bold = True; rl1.font.size = WPt(13); rl1.font.color.rgb = W_BLANCO
        rl2 = pl.add_run(nombre);     rl2.font.size = WPt(9);  rl2.font.color.rgb = W_BLANCO
        pr = cr.paragraphs[0]
        pr.add_run(desc).font.size = WPt(10)
        pr.add_run(f"\n\n[{stats}]").font.size = WPt(9)
        cr.paragraphs[0].runs[-1].font.bold = True
        doc.add_paragraph()

    doc.add_page_break()

    # ── 4. ARQUITECTURA TÉCNICA ──────────────────────────────────────
    w_heading(doc, "4. ARQUITECTURA TÉCNICA", 1, W_AZUL_OSC)
    doc.add_paragraph(
        "La solución está construida sobre tecnología robusta, de bajo costo y mantenible "
        "por el equipo interno. No requiere servidores externos ni licencias adicionales."
    )
    doc.add_paragraph()

    tech_rows = [
        ("Lenguaje",           "Python 3",                "Alto nivel, ampliamente adoptado en analítica de datos"),
        ("Automatización Web", "Playwright (Microsoft)",   "Controla Chrome para interactuar con el WMS como un usuario real"),
        ("Credenciales",       "python-dotenv / .env",    "Contraseñas almacenadas fuera del código fuente"),
        ("Almacenamiento",     "OneDrive for Business",   "Reportes sincronizados automáticamente a SharePoint"),
        ("Ejecución",          "Windows Task Scheduler",  "Lanza los 3 módulos en secuencia L-V a las 8:00 AM"),
        ("Modo silencioso",    "headless=True",            "El navegador opera en segundo plano sin ventanas visibles"),
        ("Monitoreo",          "Logs + Correo Gmail",      "Registro completo + notificación por email al terminar"),
        ("Dashboards",         "Power BI / Power Query",  "Consume archivos con nombres fijos desde OneDrive/SharePoint"),
    ]

    tech_t = doc.add_table(rows=len(tech_rows) + 1, cols=3)
    tech_t.style = "Table Grid"
    for i, h_t in enumerate(["Componente", "Tecnología", "Descripción"]):
        c = tech_t.cell(0, i); w_cell_bg(c, "1F4E79")
        rr = c.paragraphs[0].add_run(h_t)
        rr.font.color.rgb = W_BLANCO; rr.font.bold = True; rr.font.size = WPt(10)
        c.paragraphs[0].alignment = WD_ALIGN_PARAGRAPH.CENTER

    for ri, (comp, tec, desc) in enumerate(tech_rows, 1):
        bg = "DEEAF1" if ri % 2 == 0 else "FFFFFF"
        for ci, txt in enumerate([comp, tec, desc]):
            c = tech_t.cell(ri, ci); w_cell_bg(c, bg)
            rr = c.paragraphs[0].add_run(txt); rr.font.size = WPt(9)
            if ci < 2: rr.font.bold = True

    doc.add_page_break()

    # ── 5. RESULTADOS E IMPACTO ──────────────────────────────────────
    w_heading(doc, "5. RESULTADOS E IMPACTO EN EL NEGOCIO", 1, W_AZUL_OSC)
    doc.add_paragraph()

    comp_data = [
        ("INDICADOR",                "ANTES",                    "DESPUÉS"),
        ("Tiempo de descarga",       "~2 horas/día manual",      "~20 min automático"),
        ("Disponibilidad del dato",  "9:30 – 10:00 AM",          "8:20 – 8:30 AM"),
        ("Consistencia de archivos", "Variable",                  "100% estandarizado"),
        ("Riesgo de error humano",   "Alto",                      "Eliminado"),
        ("Dependencia de persona",   "Alta",                      "Nula"),
        ("Ventana de intervención",  "Ventanas visibles (lento)", "Ejecución silenciosa en 2.º plano"),
    ]

    comp_t = doc.add_table(rows=len(comp_data), cols=3)
    comp_t.style = "Table Grid"
    for ri, (c1, c2, c3) in enumerate(comp_data):
        is_h = (ri == 0)
        for ci, txt in enumerate([c1, c2, c3]):
            cell = comp_t.cell(ri, ci)
            if is_h:
                bg = "1F4E79"; fc = W_BLANCO; bld = True
            elif ci == 0:
                bg = "DEEAF1"; fc = W_AZUL_OSC; bld = True
            elif ci == 1:
                bg = "FCE4D6"; fc = W_ROJO; bld = False
            else:
                bg = "E2EFDA"; fc = WRGBColor(0x37, 0x5C, 0x23); bld = False
            w_cell_bg(cell, bg)
            cell.paragraphs[0].alignment = WD_ALIGN_PARAGRAPH.CENTER
            rr = cell.paragraphs[0].add_run(txt)
            rr.font.color.rgb = fc; rr.font.bold = bld; rr.font.size = WPt(10)

    doc.add_paragraph()
    w_caja(doc,
        "IMPACTO ANUAL ESTIMADO: Con ~2 horas de trabajo calificado liberadas por día hábil "
        "(250 días/año), el proyecto equivale a recuperar más de 500 horas de trabajo al año. "
        "Ese tiempo se reinvierte en análisis de valor, mejoras de proceso y soporte estratégico.",
        bg_hex="E2EFDA",
        text_color=WRGBColor(0x37, 0x5C, 0x23),
        bold=True)

    doc.add_page_break()

    # ── 6. BENEFICIOS ESTRATÉGICOS ───────────────────────────────────
    w_heading(doc, "6. BENEFICIOS ESTRATÉGICOS", 1, W_AZUL_OSC)

    beneficios_w = [
        ("Continuidad operacional",
         "El sistema opera independientemente de vacaciones, ausencias o carga de trabajo. "
         "Los reportes estarán disponibles en OneDrive/SharePoint a las 8:30 AM todos los días hábiles."),
        ("Calidad del dato garantizada",
         "Nombres de archivo fijos y estructura de carpetas estandarizada eliminan la posibilidad "
         "de que Power BI falle por inconsistencias. Cada CSV descargado es validado estructuralmente."),
        ("Trazabilidad y auditoría",
         "Cada ejecución genera un log con timestamp, cliente, estado (OK/error) y ruta del archivo. "
         "Adicionalmente, se envía un correo resumen al equipo al finalizar."),
        ("Escalabilidad",
         "Agregar un nuevo cliente o depósito requiere modificar una sola línea en el script. "
         "No se necesita conocimiento técnico avanzado para el mantenimiento."),
        ("Base para Analytics Avanzado",
         "Con datos consistentes fluyendo diariamente a SharePoint, el equipo puede construir "
         "dashboards en Power BI, análisis de tendencias y alertas automáticas de forma sostenible."),
    ]

    for tit_b, desc_b in beneficios_w:
        w_heading(doc, f"► {tit_b}", 2, W_AZUL_MED)
        doc.add_paragraph(desc_b)
        doc.add_paragraph()

    doc.add_page_break()

    # ── 7. PRÓXIMOS PASOS ────────────────────────────────────────────
    w_heading(doc, "7. PRÓXIMOS PASOS Y ROADMAP", 1, W_AZUL_OSC)
    doc.add_paragraph()

    roadmap_w = [
        ("Fase 6",   "✅ En curso",       "Modo headless activado — sin ventanas visibles",                    "70AD47"),
        ("Fase 5",   "⏳ Pendiente IT",   "Azure AD: correo único + SharePoint directo",                       "ED7D31"),
        ("Fase 8",   "🔮 Planificado",    "Dashboards Power BI sobre datos automatizados",                     "2E75B6"),
        ("Futuro",   "🔍 Exploración",    "Alertas automáticas ante anomalías de stock o staging",             "595959"),
    ]

    road_t = doc.add_table(rows=len(roadmap_w) + 1, cols=3)
    road_t.style = "Table Grid"
    for i, h_t in enumerate(["Fase", "Estado", "Descripción"]):
        c = road_t.cell(0, i); w_cell_bg(c, "1F4E79")
        rr = c.paragraphs[0].add_run(h_t)
        rr.font.color.rgb = W_BLANCO; rr.font.bold = True; rr.font.size = WPt(10)

    for ri, (fase, estado, desc, color) in enumerate(roadmap_w, 1):
        bg = "F2F2F2" if ri % 2 == 0 else "FFFFFF"
        for ci, txt in enumerate([fase, estado, desc]):
            c = road_t.cell(ri, ci)
            w_cell_bg(c, color + "30" if ci == 1 else bg)
            rr = c.paragraphs[0].add_run(txt); rr.font.size = WPt(10)
            if ci < 2: rr.font.bold = True

    doc.add_paragraph()
    w_caja(doc,
        "ACCIÓN REQUERIDA: Para habilitar el correo único y el envío directo a SharePoint se "
        "requiere aprobación IT para crear App Registration en Azure AD. "
        "Contacto: José Contreras — jcontreras@tinetservices.cl",
        bg_hex="FFF2CC")

    doc.add_page_break()

    # ── 8. CONCLUSIÓN ────────────────────────────────────────────────
    w_heading(doc, "8. CONCLUSIÓN", 1, W_AZUL_OSC)
    doc.add_paragraph(
        "El proyecto de Automatización de Reportes WMS representa un salto cualitativo en la "
        "madurez digital del área de Control de Gestión de Egakat SPA. No se trata únicamente "
        "de ahorrar tiempo: se trata de construir una base de datos confiable, consistente y "
        "disponible que permita tomar decisiones más rápidas y más informadas."
    )
    doc.add_paragraph()
    doc.add_paragraph(
        "Hoy, los tres módulos están en producción, ejecutándose de forma autónoma cada mañana. "
        "Los datos fluyen desde el WMS hasta OneDrive/SharePoint sin intervención humana. "
        "El equipo puede enfocarse en analizar la información, no en recopilarla."
    )
    doc.add_paragraph()
    w_caja(doc,
        "\"De recopilar datos a generar valor — ese es el cambio fundamental que este proyecto "
        "trae al área de Control de Gestión de Egakat SPA.\"",
        bg_hex="1F4E79",
        text_color=W_BLANCO,
        bold=True,
        font_size=12)

    # ── Guardar ──────────────────────────────────────────────────────
    ruta_w = os.path.join(CARPETA, "Informe_Automatizacion_WMS_Egakat.docx")
    doc.save(ruta_w)
    print(f"  ✅ Word guardado: {ruta_w}")
    return ruta_w


# ════════════════════════════════════════════════════════════════════════
#  GENERAR POWERPOINT
# ════════════════════════════════════════════════════════════════════════

def generar_ppt():
    prs = Presentation()
    prs.slide_width  = PInches(13.33)
    prs.slide_height = PInches(7.5)
    BLANK = prs.slide_layouts[6]

    # ── SLIDE 1: Portada ─────────────────────────────────────────────
    s1 = prs.slides.add_slide(BLANK)
    p_bg(s1, P_AZUL_OSC)
    p_rect(s1, 0, 0, 0.35, 7.5, P_VERDE)
    p_txt(s1, "AUTOMATIZACIÓN DE REPORTES WMS",
          0.7, 1.9, 12.0, 1.4, size=34, bold=True, color=P_BLANCO)
    p_txt(s1, "Transformación Digital del Control de Gestión — Egakat SPA",
          0.7, 3.4, 11.5, 0.8, size=17, color=P_VERDE)
    p_txt(s1, f"Sócrates Cabral   |   Head of Control Management & CI   |   {HOY}",
          0.7, 6.6, 12.0, 0.6, size=11, color=P_GRIS)

    # ── SLIDE 2: El Problema ─────────────────────────────────────────
    s2 = prs.slides.add_slide(BLANK)
    p_bg(s2, P_GRIS_CLR)
    p_rect(s2, 0, 0, 13.33, 1.2, P_AZUL_OSC)
    p_txt(s2, "EL PROBLEMA: ANTES DE LA AUTOMATIZACIÓN",
          0.3, 0.18, 12.5, 0.9, size=22, bold=True, color=P_BLANCO)

    probs = [
        ("~2 HORAS\nDIARIAS",     "Trabajo manual y repetitivo\npara descargar 27 reportes"),
        ("RIESGO DE\nERROR",      "Archivos mal nombrados\nrompen los dashboards Power BI"),
        ("DEPENDENCIA\nPERSONAL", "Si falta alguien,\nlos datos no llegan a tiempo"),
        ("DATO\nTARDIO",          "Información disponible\na media mañana"),
    ]
    colores_p = [PRGBColor(0xC0,0x00,0x00), PRGBColor(0xED,0x7D,0x31),
                 PRGBColor(0x1F,0x4E,0x79), PRGBColor(0x59,0x59,0x59)]
    for i, ((tit, desc), col) in enumerate(zip(probs, colores_p)):
        x = 0.4 + i * 3.2
        p_rect(s2, x, 1.5, 2.9, 4.6, PRGBColor(0xFF,0xFF,0xFF))
        p_rect(s2, x, 1.5, 2.9, 0.5, col)
        p_txt(s2, tit, x+0.1, 2.1, 2.7, 1.0, size=14, bold=True, color=col, align=PP_ALIGN.CENTER)
        p_txt(s2, desc, x+0.1, 3.3, 2.7, 1.8, size=11, color=P_NEGRO, align=PP_ALIGN.CENTER)

    # ── SLIDE 3: La Solución ─────────────────────────────────────────
    s3 = prs.slides.add_slide(BLANK)
    p_bg(s3, P_GRIS_CLR)
    p_rect(s3, 0, 0, 13.33, 1.2, P_AZUL_MED)
    p_txt(s3, "LA SOLUCIÓN: 3 MÓDULOS DE AUTOMATIZACIÓN",
          0.3, 0.18, 12.5, 0.9, size=22, bold=True, color=P_BLANCO)
    p_txt(s3,
          "Python + Playwright ejecuta automáticamente L-V a las 8:00 AM: "
          "27 reportes descargados y disponibles en OneDrive sin intervención humana.",
          0.5, 1.3, 12.3, 0.8, size=13, color=P_NEGRO)

    mods_s3 = [
        (P_AZUL_MED, "MÓDULO 1", "Stock WMS\nSemanal",        "3 reportes\n3 centros de distribución\n~5 minutos"),
        (P_NARANJA,  "MÓDULO 2", "Staging\nIN / OUT",          "16 reportes\n16 empresas clientes\n~10 minutos"),
        (P_VERDE,    "MÓDULO 3", "Consulta de\nPosiciones",    "8 reportes\n4 depósitos\n~4 minutos"),
    ]
    for i, (col, cod, nom, stats) in enumerate(mods_s3):
        x = 0.8 + i * 4.0
        p_rect(s3, x, 2.3, 3.6, 0.9, col)
        p_txt(s3, f"{cod}: {nom}", x+0.1, 2.35, 3.4, 0.85, size=13, bold=True, color=P_BLANCO, align=PP_ALIGN.CENTER)
        p_rect(s3, x, 3.2, 3.6, 2.6, PRGBColor(0xFF,0xFF,0xFF))
        p_txt(s3, stats, x+0.1, 3.3, 3.4, 2.4, size=13, color=P_NEGRO, align=PP_ALIGN.CENTER)

    p_rect(s3, 0.5, 6.2, 12.3, 0.9, PRGBColor(0xDE,0xEA,0xF1))
    p_txt(s3, "Total: 27 reportes automatizados   |   Ejecución L-V 8:00 AM   |   Sin intervención humana",
          0.6, 6.25, 12.0, 0.75, size=12, bold=True, color=P_AZUL_OSC, align=PP_ALIGN.CENTER)

    # ── SLIDE 4: Módulo 1 ────────────────────────────────────────────
    s4 = prs.slides.add_slide(BLANK)
    p_bg(s4, P_GRIS_CLR)
    p_rect(s4, 0, 0, 13.33, 1.2, P_AZUL_MED)
    p_txt(s4, "MÓDULO 1 — STOCK WMS SEMANAL",
          0.3, 0.18, 12.5, 0.9, size=22, bold=True, color=P_BLANCO)
    p_txt(s4, "Descarga diaria del inventario completo de los 3 centros de distribución activos.",
          0.5, 1.3, 12.0, 0.6, size=13, color=P_NEGRO)
    for i, (cd, desc) in enumerate([
        ("Quilicura",         "Inventario completo CD principal"),
        ("Pudahuel",          "Inventario CD Moderno + Unitario"),
        ("Pudahuel Unitario", "Inventario unitarizado"),
    ]):
        y = 2.1 + i * 1.1
        p_rect(s4, 0.5, y, 0.6, 0.85, P_AZUL_MED)
        p_txt(s4, str(i+1), 0.5, y+0.1, 0.6, 0.65, size=18, bold=True, color=P_BLANCO, align=PP_ALIGN.CENTER)
        p_rect(s4, 1.2, y, 11.3, 0.85, PRGBColor(0xFF,0xFF,0xFF))
        p_txt(s4, f"{cd}  —  {desc}", 1.35, y+0.15, 11.0, 0.6, size=12, color=P_NEGRO)

    p_rect(s4, 0.5, 5.5, 12.0, 1.1, PRGBColor(0xDE,0xEA,0xF1))
    p_txt(s4,
          "Resultado: 3 archivos Excel con nombre estandarizado en OneDrive → "
          "Power BI los consume sin configuración adicional.",
          0.7, 5.6, 11.6, 0.9, size=12, color=P_AZUL_OSC)

    # ── SLIDE 5: Módulo 2 ────────────────────────────────────────────
    s5 = prs.slides.add_slide(BLANK)
    p_bg(s5, P_GRIS_CLR)
    p_rect(s5, 0, 0, 13.33, 1.2, P_NARANJA)
    p_txt(s5, "MÓDULO 2 — STAGING IN / OUT",
          0.3, 0.18, 12.5, 0.9, size=22, bold=True, color=P_BLANCO)
    p_txt(s5, "Descarga el reporte 'Consulta Stock con Staging In y Out' para 16 empresas clientes.",
          0.5, 1.3, 12.0, 0.6, size=13, color=P_NEGRO)

    sesiones = [
        (P_AZUL_MED, "QUILICURA",         "Cervecería ABI\nDaikin · Daikin Clientes\nDerco · Mascotas Latinas\nPochteca\n(6 clientes)"),
        (P_VERDE,    "PUDAHUEL",           "Barentz · Buraschi · Cepas Chile\nCollico · Delibest · Intime\nNativo Drinks · Tres Montes\nUnilever\n(9 clientes)"),
        (P_NARANJA,  "PUDAHUEL\nUNITARIO", "Runo SPA\n(1 cliente)"),
    ]
    for i, (col, dep, clts) in enumerate(sesiones):
        x = 0.4 + i * 4.3
        p_rect(s5, x, 2.0, 4.0, 0.8, col)
        p_txt(s5, dep, x+0.1, 2.05, 3.8, 0.75, size=13, bold=True, color=P_BLANCO, align=PP_ALIGN.CENTER)
        p_rect(s5, x, 2.8, 4.0, 3.3, PRGBColor(0xFF,0xFF,0xFF))
        p_txt(s5, clts, x+0.15, 2.9, 3.7, 3.1, size=10, color=P_NEGRO)

    p_rect(s5, 0.4, 6.3, 12.5, 0.85, PRGBColor(0xFD,0xF2,0xD0))
    p_txt(s5, "Cada cliente recibe su CSV con nombre original del WMS — trazabilidad completa.",
          0.6, 6.4, 12.0, 0.65, size=11, bold=True, color=P_NARANJA)

    # ── SLIDE 6: Módulo 3 ────────────────────────────────────────────
    s6 = prs.slides.add_slide(BLANK)
    p_bg(s6, P_GRIS_CLR)
    p_rect(s6, 0, 0, 13.33, 1.2, P_VERDE)
    p_txt(s6, "MÓDULO 3 — CONSULTA DE POSICIONES",
          0.3, 0.18, 12.5, 0.9, size=22, bold=True, color=P_BLANCO)
    p_txt(s6, "Descarga 8 reportes de posiciones (Ocupadas / Libres) para los 4 depósitos.",
          0.5, 1.3, 12.0, 0.6, size=13, color=P_NEGRO)

    pos_rows = [
        ("Quilicura",          "Posiciones Ocupadas.xlsx",            "Posiciones Libres.xlsx"),
        ("Pudahuel",           "Posiciones Ocupadas Moderno.xlsx",    "Posiciones Libres Moderno.xlsx"),
        ("Pudahuel Unitario",  "Posiciones Ocupadas Unitario.xlsx",   "Posiciones Libres Unitario.xlsx"),
        ("Pudahuel Refrigerado","Posiciones Ocupadas Refrigerado.xlsx","Posiciones Libres Refrigerado.xlsx"),
    ]
    for i, (cd, ocu, lib) in enumerate(pos_rows):
        y = 2.0 + i * 1.1
        p_rect(s6, 0.4, y, 3.3, 0.85, PRGBColor(0xFF,0xFF,0xFF))
        p_txt(s6, cd, 0.5, y+0.12, 3.1, 0.65, size=11, bold=True, color=P_VERDE)
        p_rect(s6, 3.9, y, 4.2, 0.85, PRGBColor(0xE2,0xEF,0xDA))
        p_txt(s6, ocu, 4.0, y+0.12, 4.0, 0.65, size=9, color=PRGBColor(0x37,0x5C,0x23))
        p_rect(s6, 8.3, y, 4.7, 0.85, PRGBColor(0xDE,0xEA,0xF1))
        p_txt(s6, lib, 8.4, y+0.12, 4.5, 0.65, size=9, color=P_AZUL_OSC)

    p_rect(s6, 0.4, 6.5, 12.5, 0.75, PRGBColor(0xE2,0xEF,0xDA))
    p_txt(s6, "Nombres de archivo FIJOS — Power Query los conecta directamente sin re-configurar.",
          0.6, 6.55, 12.0, 0.6, size=11, bold=True, color=PRGBColor(0x37,0x5C,0x23))

    # ── SLIDE 7: Beneficios ──────────────────────────────────────────
    s7 = prs.slides.add_slide(BLANK)
    p_bg(s7, P_AZUL_OSC)
    p_txt(s7, "IMPACTO Y BENEFICIOS",
          0.5, 0.3, 12.3, 1.0, size=30, bold=True, color=P_BLANCO)

    metricas = [
        (P_VERDE,                     "500+\nhoras/año",    "liberadas para análisis\nde valor estratégico"),
        (P_NARANJA,                   "100%\nconsistente",  "datos estandarizados\npara Power BI"),
        (P_AZUL_MED,                  "8:20 AM",            "datos disponibles antes\nde iniciar el día"),
        (PRGBColor(0xA0,0x00,0x00),   "0 errores\nhumanos", "proceso automático\nsin intervención"),
    ]
    for i, (col, metric, desc) in enumerate(metricas):
        x = 0.4 + i * 3.1
        p_rect(s7, x, 1.8, 2.9, 1.2, col)
        p_txt(s7, metric, x+0.1, 1.85, 2.7, 1.1, size=18, bold=True, color=P_BLANCO, align=PP_ALIGN.CENTER)
        p_txt(s7, desc,   x+0.1, 3.1,  2.7, 0.9, size=11, color=P_GRIS,  align=PP_ALIGN.CENTER)

    p_rect(s7, 0.5, 4.3, 12.3, 2.3, P_AZUL_MED)
    p_txt(s7,
          "\"De recopilar datos a generar valor.\"\n\n"
          "El equipo ya no dedica 2 horas diarias a descargar reportes. "
          "Ese tiempo se reinvierte en análisis, mejoras de proceso y soporte estratégico a la operación.",
          0.8, 4.4, 11.7, 2.1, size=13, italic=True, color=P_BLANCO)

    # ── SLIDE 8: Cómo funciona ───────────────────────────────────────
    s8 = prs.slides.add_slide(BLANK)
    p_bg(s8, P_GRIS_CLR)
    p_rect(s8, 0, 0, 13.33, 1.2, P_AZUL_OSC)
    p_txt(s8, "¿CÓMO FUNCIONA? — ARQUITECTURA SIMPLIFICADA",
          0.3, 0.18, 12.5, 0.9, size=22, bold=True, color=P_BLANCO)

    pasos = [
        (P_AZUL_MED, "8:00 AM",    "Programador\nde Tareas",   "Activa el proceso\nautomáticamente"),
        (P_AZUL_MED, "8:00–8:20",  "Python +\nPlaywright",     "Navega el WMS y\ndescarga 27 reportes"),
        (P_VERDE,    "8:20 AM",    "OneDrive\nfor Business",   "Archivos sincronizados\na SharePoint"),
        (P_VERDE,    "8:25 AM",    "Power BI\nDashboards",     "Reportes actualizados\nautomáticamente"),
    ]
    for i, (col, hora, comp, desc) in enumerate(pasos):
        x = 0.5 + i * 3.1
        p_rect(s8, x, 1.5, 2.8, 0.7, col)
        p_txt(s8, hora, x+0.1, 1.55, 2.6, 0.6, size=13, bold=True, color=P_BLANCO, align=PP_ALIGN.CENTER)
        p_rect(s8, x, 2.2, 2.8, 2.1, PRGBColor(0xFF,0xFF,0xFF))
        p_txt(s8, comp, x+0.1, 2.3, 2.6, 0.9, size=13, bold=True, color=col, align=PP_ALIGN.CENTER)
        p_txt(s8, desc, x+0.1, 3.2, 2.6, 1.0, size=10, color=P_NEGRO, align=PP_ALIGN.CENTER)
        if i < 3:
            p_txt(s8, "→", x+2.7, 2.4, 0.4, 0.5, size=20, bold=True, color=P_AZUL_OSC, align=PP_ALIGN.CENTER)

    p_rect(s8, 0.5, 4.6, 12.3, 1.8, PRGBColor(0xDE,0xEA,0xF1))
    p_txt(s8,
          "Tecnologías: Python 3  |  Playwright (Microsoft)  |  Windows Task Scheduler  |  "
          "OneDrive for Business  |  Power BI\n\n"
          "Modo: headless=True — el navegador opera en segundo plano sin interferir en el trabajo diario.",
          0.7, 4.7, 11.9, 1.6, size=11, color=P_AZUL_OSC)

    # ── SLIDE 9: Próximos Pasos ──────────────────────────────────────
    s9 = prs.slides.add_slide(BLANK)
    p_bg(s9, P_GRIS_CLR)
    p_rect(s9, 0, 0, 13.33, 1.2, P_AZUL_MED)
    p_txt(s9, "PRÓXIMOS PASOS — ROADMAP",
          0.3, 0.18, 12.5, 0.9, size=22, bold=True, color=P_BLANCO)

    road = [
        (P_VERDE,   "✅ En curso",      "Modo silencioso activo",     "headless=True en los 3 módulos — sin ventanas visibles"),
        (P_NARANJA, "⏳ Pendiente IT",  "Azure AD + Correo único",    "Aprobación IT para App Registration — jcontreras@tinetservices.cl"),
        (P_AZUL_MED,"🔮 Planificado",   "Power BI Dashboards",        "Dashboards operacionales sobre los datos automatizados"),
        (P_GRIS,    "🔍 Exploración",   "Alertas inteligentes",       "Notificaciones ante anomalías de stock o staging"),
    ]
    for i, (col, est, tit_r, desc_r) in enumerate(road):
        y = 1.5 + i * 1.3
        p_rect(s9, 0.4, y, 2.4, 1.0, col)
        p_txt(s9, est, 0.5, y+0.15, 2.2, 0.75, size=11, bold=True, color=P_BLANCO, align=PP_ALIGN.CENTER)
        p_rect(s9, 2.9, y, 9.8, 1.0, PRGBColor(0xFF,0xFF,0xFF))
        p_txt(s9, f"{tit_r}\n{desc_r}", 3.0, y+0.07, 9.5, 0.88, size=11, color=P_NEGRO)

    p_rect(s9, 0.4, 6.8, 12.5, 0.55, PRGBColor(0xFD,0xF2,0xD0))
    p_txt(s9, "Para avanzar en Azure AD: enviar aprobación a José Contreras (IT) — jcontreras@tinetservices.cl",
          0.6, 6.83, 12.0, 0.45, size=10, bold=True, color=P_NARANJA)

    # ── SLIDE 10: Cierre ─────────────────────────────────────────────
    s10 = prs.slides.add_slide(BLANK)
    p_bg(s10, P_AZUL_OSC)
    p_rect(s10, 0, 0, 0.35, 7.5, P_VERDE)
    p_txt(s10, "GRACIAS",
          0.7, 1.4, 12.0, 1.6, size=52, bold=True, color=P_BLANCO)
    p_txt(s10,
          "El futuro de la operación logística no está en hacer más de lo mismo,\n"
          "sino en hacer que lo repetitivo se haga solo.",
          0.7, 3.2, 12.0, 1.2, size=16, italic=True, color=P_VERDE)
    p_txt(s10,
          "Sócrates Cabral   |   Head of Control Management & Continuous Improvement\n"
          "Egakat SPA   |   socrates.cabral@egakat.cl",
          0.7, 5.6, 12.0, 1.0, size=12, color=P_GRIS)

    # ── Guardar ──────────────────────────────────────────────────────
    ruta_p = os.path.join(CARPETA, "Presentacion_Automatizacion_WMS_Egakat.pptx")
    prs.save(ruta_p)
    print(f"  ✅ PowerPoint guardado: {ruta_p}")
    return ruta_p


# ════════════════════════════════════════════════════════════════════════
#  MAIN
# ════════════════════════════════════════════════════════════════════════

def main():
    print("=" * 60)
    print("  Generando documentos — WMS Egakat Automatización")
    print("=" * 60)
    print()
    ruta_w = generar_word()
    ruta_p = generar_ppt()
    print()
    print("=" * 60)
    print("  Documentos generados correctamente:")
    print(f"  Word : {ruta_w}")
    print(f"  PPT  : {ruta_p}")
    print("=" * 60)


if __name__ == "__main__":
    main()
