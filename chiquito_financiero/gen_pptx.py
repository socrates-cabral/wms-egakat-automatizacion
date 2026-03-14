import sys
sys.stdout.reconfigure(encoding="utf-8")
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# ── Paleta de colores ──────────────────────────────────────────────────────────
AZUL       = RGBColor(0x1F, 0x6F, 0xEB)
AZUL_OSCURO= RGBColor(0x0D, 0x47, 0xA1)
VERDE      = RGBColor(0x1E, 0x8B, 0x4C)
ROJO       = RGBColor(0xC0, 0x39, 0x2B)
AMBAR      = RGBColor(0x9A, 0x7D, 0x0A)
GRIS       = RGBColor(0x5D, 0x6D, 0x7E)
GRIS_CLARO = RGBColor(0xF2, 0xF3, 0xF4)
BLANCO     = RGBColor(0xFF, 0xFF, 0xFF)
NEGRO      = RGBColor(0x1A, 0x1A, 0x2E)
FONDO      = RGBColor(0x0D, 0x11, 0x17)
PANEL      = RGBColor(0x16, 0x1B, 0x22)

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

BLANK = prs.slide_layouts[6]  # Blank layout

# ── Helpers ────────────────────────────────────────────────────────────────────
def add_slide():
    return prs.slides.add_slide(BLANK)

def rect(slide, l, t, w, h, fill_rgb=None, line_rgb=None, line_width=None):
    from pptx.util import Pt
    shape = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    shape.line.fill.background()
    if fill_rgb:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_rgb
    else:
        shape.fill.background()
    if line_rgb:
        shape.line.color.rgb = line_rgb
        if line_width:
            shape.line.width = Pt(line_width)
    else:
        shape.line.fill.background()
    return shape

def textbox(slide, text, l, t, w, h, size=18, bold=False, color=NEGRO, align=PP_ALIGN.LEFT,
            italic=False, wrap=True):
    txBox = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txBox

def add_table(slide, headers, rows, l, t, w, h, header_fill=AZUL, alt_fill=GRIS_CLARO):
    from pptx.util import Pt
    n_cols = len(headers)
    n_rows = len(rows) + 1
    tbl = slide.shapes.add_table(n_rows, n_cols, Inches(l), Inches(t), Inches(w), Inches(h)).table
    col_w = Inches(w) // n_cols
    for i in range(n_cols):
        tbl.columns[i].width = col_w
    # Header
    for ci, h_text in enumerate(headers):
        cell = tbl.cell(0, ci)
        cell.text = h_text
        cell.fill.solid()
        cell.fill.fore_color.rgb = header_fill
        for para in cell.text_frame.paragraphs:
            for run in para.runs:
                run.font.bold = True
                run.font.color.rgb = BLANCO
                run.font.size = Pt(11)
            para.alignment = PP_ALIGN.CENTER
    # Rows
    for ri, row in enumerate(rows):
        for ci, val in enumerate(row):
            cell = tbl.cell(ri + 1, ci)
            cell.text = str(val)
            if ri % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = alt_fill
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = BLANCO
            for para in cell.text_frame.paragraphs:
                for run in para.runs:
                    run.font.size = Pt(10)
    return tbl

def kpi_card(slide, l, t, valor, label, color=AZUL, ancho=2.5, alto=1.1):
    r = rect(slide, l, t, ancho, alto, fill_rgb=PANEL, line_rgb=color, line_width=1.5)
    textbox(slide, valor, l+0.1, t+0.08, ancho-0.2, 0.55, size=22, bold=True, color=color, align=PP_ALIGN.CENTER)
    textbox(slide, label, l+0.1, t+0.62, ancho-0.2, 0.35, size=9, color=GRIS, align=PP_ALIGN.CENTER)

def slide_bg(slide, color=FONDO):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def header_bar(slide, titulo, subtitulo=""):
    rect(slide, 0, 0, 13.33, 1.1, fill_rgb=PANEL)
    rect(slide, 0, 1.0, 13.33, 0.05, fill_rgb=AZUL)
    textbox(slide, titulo, 0.3, 0.1, 9, 0.55, size=24, bold=True, color=BLANCO)
    if subtitulo:
        textbox(slide, subtitulo, 0.3, 0.65, 9, 0.35, size=13, color=GRIS, italic=True)
    textbox(slide, "Chiquito Finanzas | Mar 2026", 10.5, 0.1, 2.7, 0.4, size=9, color=GRIS, align=PP_ALIGN.RIGHT)

def footer(slide):
    rect(slide, 0, 7.2, 13.33, 0.3, fill_rgb=PANEL)
    textbox(slide, "Chiquito Finanzas  |  Diagnóstico Financiero  |  Marzo 2026  |  Confidencial", 0.3, 7.22, 12, 0.25, size=8, color=GRIS)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — PORTADA
# ══════════════════════════════════════════════════════════════════════════════
slide = add_slide()
slide_bg(slide, FONDO)
rect(slide, 0, 0, 13.33, 7.5, fill_rgb=FONDO)
# Línea superior
rect(slide, 0, 0, 13.33, 0.08, fill_rgb=AZUL)
# Panel central
rect(slide, 1.5, 1.8, 10.3, 4.0, fill_rgb=PANEL, line_rgb=AZUL, line_width=1)
rect(slide, 1.5, 1.8, 10.3, 0.08, fill_rgb=AZUL)

textbox(slide, "DIAGNÓSTICO FINANCIERO", 2.0, 2.1, 9.5, 0.7, size=13, bold=True, color=GRIS, align=PP_ALIGN.CENTER)
textbox(slide, "Chiquito Finanzas", 2.0, 2.6, 9.5, 1.0, size=40, bold=True, color=BLANCO, align=PP_ALIGN.CENTER)
textbox(slide, "Taller de Muebles  |  Macul, Santiago, Chile", 2.0, 3.55, 9.5, 0.5, size=16, color=GRIS, align=PP_ALIGN.CENTER, italic=True)

rect(slide, 4.5, 4.2, 4.3, 0.04, fill_rgb=AZUL)

textbox(slide, "Estado: DÉFICIT OPERATIVO — Requiere acción inmediata", 2.0, 4.35, 9.5, 0.4, size=13, bold=True, color=ROJO, align=PP_ALIGN.CENTER)
textbox(slide, "Presentado por: Sócrates Cabral  |  Control de Gestión y Mejora Continua  |  Egakat SPA", 2.0, 4.85, 9.5, 0.35, size=10, color=GRIS, align=PP_ALIGN.CENTER)
textbox(slide, "Marzo 2026", 2.0, 5.2, 9.5, 0.35, size=12, color=GRIS, align=PP_ALIGN.CENTER)
# Línea inferior
rect(slide, 0, 7.42, 13.33, 0.08, fill_rgb=AZUL)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — RESUMEN EJECUTIVO
# ══════════════════════════════════════════════════════════════════════════════
slide = add_slide()
slide_bg(slide, FONDO)
header_bar(slide, "Resumen Ejecutivo", "El diagnóstico en 60 segundos")
footer(slide)

# Columna izquierda — situación
rect(slide, 0.3, 1.2, 5.8, 5.8, fill_rgb=PANEL, line_rgb=ROJO, line_width=1)
rect(slide, 0.3, 1.2, 5.8, 0.35, fill_rgb=ROJO)
textbox(slide, "⚠  SITUACIÓN ACTUAL", 0.5, 1.22, 5.4, 0.3, size=11, bold=True, color=BLANCO)

puntos_situacion = [
    "Ventas promedio: $1,842,930/mes",
    "Gastos promedio: $1,948,700/mes",
    "Resultado: -$105,770/mes (déficit)",
    "Deuda total: $26,451,837",
    "Cuotas mensuales: $918,903 (50% del ingreso)",
    "Punto de Equilibrio: $4,250,896/mes",
    "Ventas actuales = 43% del PE",
    "El negocio necesita más que duplicar sus ventas",
]
for i, p in enumerate(puntos_situacion):
    textbox(slide, "• " + p, 0.5, 1.65 + i*0.57, 5.4, 0.5, size=11, color=BLANCO if i < 3 else GRIS)

# Columna derecha — diagnóstico
rect(slide, 6.5, 1.2, 6.5, 5.8, fill_rgb=PANEL, line_rgb=AMBAR, line_width=1)
rect(slide, 6.5, 1.2, 6.5, 0.35, fill_rgb=AMBAR)
textbox(slide, "🔍  5 CAUSAS DEL DÉFICIT", 6.7, 1.22, 6.1, 0.3, size=11, bold=True, color=BLANCO)

causas = [
    ("1.", "Deuda aplastante", "Cuotas = 50% del ingreso promedio"),
    ("2.", "Alquiler excesivo", "$700K/mes = 38% del ingreso (normal: 10-15%)"),
    ("3.", "Ventas insuficientes", "Vende $1.8M vs. necesita $3.95M"),
    ("4.", "Gastos personales mezclados", "~$80K/mes de gastos privados en la caja"),
    ("5.", "TCs en mora o riesgo de mora", "Genera tasa TMC 2.75%/mes — la más alta"),
]
for i, (num, titulo, desc) in enumerate(causas):
    y = 1.65 + i * 1.0
    textbox(slide, num, 6.7, y, 0.4, 0.35, size=20, bold=True, color=AMBAR)
    textbox(slide, titulo, 7.1, y, 5.5, 0.35, size=13, bold=True, color=BLANCO)
    textbox(slide, desc, 7.1, y + 0.35, 5.5, 0.4, size=10, color=GRIS, italic=True)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — KPIs FINANCIEROS
# ══════════════════════════════════════════════════════════════════════════════
slide = add_slide()
slide_bg(slide, FONDO)
header_bar(slide, "KPIs Financieros", "Indicadores clave del desempeño operativo — Datos reales Nov 2025 a Mar 2026")
footer(slide)

# Fila 1 — KPIs principales
kpis_row1 = [
    ("$1,842,930", "Ingreso prom/mes", VERDE),
    ("$1,948,700", "Gasto prom/mes", ROJO),
    ("-$105,770", "Resultado neto prom", ROJO),
    ("43%", "% Punto de Equilibrio", ROJO),
]
for i, (val, lbl, col) in enumerate(kpis_row1):
    kpi_card(slide, 0.3 + i * 3.15, 1.25, val, lbl, col, ancho=3.0, alto=1.15)

# Fila 2 — KPIs deuda
kpis_row2 = [
    ("$26,451,837", "Deuda total consolidada", ROJO),
    ("$918,903", "Cuotas mensuales totales", AMBAR),
    ("50%", "% ingreso destinado a cuotas", AMBAR),
    ("8", "Instrumentos de deuda activos", AZUL),
]
for i, (val, lbl, col) in enumerate(kpis_row2):
    kpi_card(slide, 0.3 + i * 3.15, 2.6, val, lbl, col, ancho=3.0, alto=1.15)

# Fila 3 — Tabla flujo mensual
rect(slide, 0.3, 3.95, 12.7, 3.3, fill_rgb=PANEL)
textbox(slide, "Flujo mensual real (Nov 2025 – Mar 2026)", 0.5, 4.0, 8, 0.35, size=12, bold=True, color=BLANCO)
add_table(slide,
    ["Mes", "Ingresos", "Gastos", "Resultado", "Estado"],
    [
        ["Nov-2025", "$1,721,170", "$2,025,470", "-$304,300", "❌ Déficit"],
        ["Dic-2025", "$3,024,913", "$2,715,420", "+$309,493", "✅ Superávit"],
        ["Ene-2026", "$1,625,820", "$1,686,982", "-$61,162", "❌ Déficit"],
        ["Feb-2026", "$1,964,928", "$1,617,387", "+$347,541", "✅ Superávit"],
        ["Mar-2026*", "$2,400,000", "$2,200,000", "+$200,000", "⚠ Estimado"],
    ],
    0.3, 4.4, 12.7, 2.75
)
textbox(slide, "*Mar-2026 estimado con datos parciales al 11-Mar", 0.5, 7.1, 8, 0.2, size=8, color=GRIS, italic=True)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — ESTRUCTURA DE COSTOS Y DEUDA
# ══════════════════════════════════════════════════════════════════════════════
slide = add_slide()
slide_bg(slide, FONDO)
header_bar(slide, "Estructura de Costos y Deuda", "Dónde va cada peso que entra al negocio")
footer(slide)

# Panel costos fijos
rect(slide, 0.3, 1.2, 5.8, 3.0, fill_rgb=PANEL, line_rgb=AMBAR, line_width=1)
rect(slide, 0.3, 1.2, 5.8, 0.3, fill_rgb=AMBAR)
textbox(slide, "COSTOS FIJOS MENSUALES  — Total: $1,912,903", 0.5, 1.22, 5.4, 0.28, size=10, bold=True, color=BLANCO)

costos = [
    ("Cuotas bancarias (8 deudas)", "$918,903", "48%", ROJO),
    ("Alquiler del taller", "$700,000", "37%", AMBAR),
    ("Gasolina camión", "$100,000", "5%", GRIS),
    ("Luz, agua y servicios", "$55,000", "3%", GRIS),
    ("Otros (teléfono, internet...)", "$139,000", "7%", GRIS),
]
for i, (nombre, monto, pct, col) in enumerate(costos):
    y = 1.58 + i * 0.48
    rect(slide, 0.4, y, 0.12, 0.28, fill_rgb=col)
    textbox(slide, nombre, 0.6, y, 3.5, 0.28, size=10, color=BLANCO)
    textbox(slide, monto, 4.0, y, 1.0, 0.28, size=10, bold=True, color=col, align=PP_ALIGN.RIGHT)
    textbox(slide, pct, 5.05, y, 0.7, 0.28, size=10, color=GRIS, align=PP_ALIGN.RIGHT)

# Panel deuda por tipo
rect(slide, 6.5, 1.2, 6.5, 3.0, fill_rgb=PANEL, line_rgb=ROJO, line_width=1)
rect(slide, 6.5, 1.2, 6.5, 0.3, fill_rgb=ROJO)
textbox(slide, "DEUDA POR INSTRUMENTO  — Total: $26,451,837", 6.7, 1.22, 6.1, 0.28, size=10, bold=True, color=BLANCO)

deudas = [
    ("Crédito Foton (camión)", "$9,517,195", "36%", "$264,366/mes", "1.2%/mes"),
    ("Banco Itaú (crédito)", "$5,749,547", "22%", "$154,028/mes", "2.8%/mes"),
    ("Banco Estado (crédito)", "$5,600,000", "21%", "$174,437/mes", "3.1%/mes"),
    ("Banco Santander (TC)", "$3,760,935", "14%", "$109,000/mes", "2.8%/mes"),
    ("Líneas de crédito", "$2,360,000", "9%", "$71,660/mes", "3.1%/mes"),
    ("CMR Falabella", "$1,607,443", "6%", "$80,000/mes", "3.3%/mes"),
    ("Hermana (dólares)", "$1,050,000", "4%", "$0/mes", "0%"),
    ("Seguro Foton", "$0", "—", "$65,412/mes", "fijo"),
]
for i, (nombre, saldo, pct, cuota, tasa) in enumerate(deudas):
    y = 1.58 + i * 0.33
    tasa_val = float(tasa.replace('%/mes','').replace('fijo','0').replace('%',''))
    col_tasa = ROJO if tasa_val >= 2.8 else (AMBAR if tasa_val >= 1.2 else VERDE)
    textbox(slide, nombre, 6.6, y, 2.8, 0.3, size=9.5, color=BLANCO)
    textbox(slide, saldo, 9.4, y, 1.3, 0.3, size=9.5, bold=True, color=BLANCO, align=PP_ALIGN.RIGHT)
    textbox(slide, cuota, 10.7, y, 1.3, 0.3, size=9, color=AMBAR, align=PP_ALIGN.RIGHT)
    textbox(slide, tasa, 12.0, y, 0.9, 0.3, size=9, color=col_tasa, align=PP_ALIGN.RIGHT)

# Indicador crítico
rect(slide, 0.3, 4.35, 12.7, 0.85, fill_rgb=RGBColor(0x2D, 0x1B, 0x1B), line_rgb=ROJO, line_width=1)
textbox(slide, "⚠  El 48% de los costos fijos son cuotas bancarias. Sin reducir la deuda o aumentar drásticamente las ventas, el negocio no puede ser viable. Las tarjetas de crédito (tasa 2.8-3.3%/mes) son el foco de acción prioritario.", 0.5, 4.4, 12.3, 0.75, size=11, bold=True, color=ROJO)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — PUNTO DE EQUILIBRIO
# ══════════════════════════════════════════════════════════════════════════════
slide = add_slide()
slide_bg(slide, FONDO)
header_bar(slide, "Punto de Equilibrio", "¿Cuánto necesita vender el negocio para sobrevivir?")
footer(slide)

# Visualización del PE — barra de progreso
rect(slide, 0.3, 1.25, 12.7, 2.5, fill_rgb=PANEL)

textbox(slide, "VENTAS ACTUALES: $1,842,930/mes", 0.5, 1.3, 6, 0.4, size=13, bold=True, color=VERDE)
textbox(slide, "PUNTO DE EQUILIBRIO: $4,250,896/mes", 6.5, 1.3, 6.3, 0.4, size=13, bold=True, color=AMBAR)

# Barra de progreso visual
rect(slide, 0.5, 1.85, 12.3, 0.6, fill_rgb=RGBColor(0x21, 0x26, 0x2D))
# Progreso actual (43%)
rect(slide, 0.5, 1.85, 5.29, 0.6, fill_rgb=ROJO)
textbox(slide, "43%", 2.5, 1.88, 1.5, 0.55, size=18, bold=True, color=BLANCO, align=PP_ALIGN.CENTER)
textbox(slide, "Ventas actuales", 0.5, 2.5, 5.3, 0.3, size=9, color=GRIS, align=PP_ALIGN.CENTER)
textbox(slide, "Punto de Equilibrio (100%)", 5.5, 2.5, 5.3, 0.3, size=9, color=GRIS, align=PP_ALIGN.CENTER)
# Línea PE
rect(slide, 12.3, 1.75, 0.03, 0.8, fill_rgb=AMBAR)
textbox(slide, "▼ PE", 12.0, 1.7, 0.6, 0.3, size=9, color=AMBAR)
# Meta 70%
rect(slide, 9.11, 1.75, 0.03, 0.8, fill_rgb=VERDE)
textbox(slide, "Meta mínima 70%", 8.5, 1.7, 1.5, 0.3, size=8, color=VERDE, align=PP_ALIGN.CENTER)

# Tabla de escenarios
rect(slide, 0.3, 3.85, 12.7, 3.35, fill_rgb=PANEL)
textbox(slide, "¿Qué necesita cambiar para alcanzar el equilibrio?", 0.5, 3.9, 10, 0.35, size=13, bold=True, color=BLANCO)
add_table(slide,
    ["Escenario", "Ventas necesarias", "Cambios requeridos", "% PE", "Factibilidad"],
    [
        ["Actual (sin cambios)", "$4,250,896", "Ninguno — solo subir ventas", "43% → 100%", "Muy difícil"],
        ["Renegociar alquiler", "$3,582,484", "Alquiler $700K → $400K", "43% → 100%", "Difícil pero posible"],
        ["Renegociar + subir precio", "$3,186,484", "Alquiler $400K + precios +15%", "43% → 100%", "Alcanzable en 6-12 meses"],
        ["Inyección + renegociar", "$2,962,484", "Crédito BCI + alquiler $400K", "43% → 100%", "Alcanzable en 3-6 meses"],
        ["Escenario óptimo 2027", "$2,500,000", "Todas las acciones ejecutadas", "74%", "Meta realista a 18 meses"],
    ],
    0.3, 4.3, 12.7, 2.8
)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — SIMULACIONES
# ══════════════════════════════════════════════════════════════════════════════
slide = add_slide()
slide_bg(slide, FONDO)
header_bar(slide, "Simulación de Escenarios", "¿Qué pasa si...? — Cuatro caminos posibles")
footer(slide)

escenarios = [
    ("📍 ACTUAL", "$1,860,000", "$700,000", "$363,097", "-$540,000/mes", ROJO, "Insostenible — pérdidas continuas"),
    ("🌱 OPTIMISTA", "$3,000,000", "$700,000", "$363,097", "+$76,500/mes", VERDE, "Viable solo si ventas triplican — difícil"),
    ("🤝 RENEGOCIADO", "$2,200,000", "$400,000", "$200,000", "+$149,000/mes", VERDE, "Alcanzable con negociaciones urgentes"),
    ("⚖ EQUILIBRIO", "$3,950,000", "$700,000", "$363,097", "$0/mes", AMBAR, "Requiere doblar ventas — meta 2027"),
]

for i, (nombre, ventas, alquiler, cuotas, resultado, col, nota) in enumerate(escenarios):
    x = 0.3 + (i % 2) * 6.55
    y = 1.25 + (i // 2) * 2.9
    rect(slide, x, y, 6.2, 2.65, fill_rgb=PANEL, line_rgb=col, line_width=1.5)
    rect(slide, x, y, 6.2, 0.38, fill_rgb=col)
    textbox(slide, nombre, x + 0.15, y + 0.05, 5.9, 0.3, size=13, bold=True, color=BLANCO)
    # Datos
    datos = [("Ventas/mes", ventas), ("Alquiler", alquiler), ("Cuotas TC", cuotas)]
    for j, (lbl, val) in enumerate(datos):
        textbox(slide, lbl + ":", x + 0.2, y + 0.5 + j*0.38, 2.5, 0.35, size=10, color=GRIS)
        textbox(slide, val, x + 2.5, y + 0.5 + j*0.38, 3.5, 0.35, size=10, bold=True, color=BLANCO, align=PP_ALIGN.RIGHT)
    # Resultado
    rect(slide, x + 0.2, y + 1.7, 5.8, 0.38, fill_rgb=RGBColor(0x21, 0x26, 0x2D))
    textbox(slide, "Resultado: " + resultado, x + 0.3, y + 1.72, 5.6, 0.35, size=13, bold=True, color=col, align=PP_ALIGN.CENTER)
    textbox(slide, nota, x + 0.2, y + 2.15, 5.8, 0.35, size=9, color=GRIS, italic=True)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — INYECCIÓN DE CAPITAL
# ══════════════════════════════════════════════════════════════════════════════
slide = add_slide()
slide_bg(slide, FONDO)
header_bar(slide, "Opción: Inyección de Capital", "Crédito BCI $10M + Aporte familiar $2.2M para refinanciar deuda costosa")
footer(slide)

# Izquierda — propuesta
rect(slide, 0.3, 1.2, 5.8, 5.8, fill_rgb=PANEL, line_rgb=AZUL, line_width=1)
rect(slide, 0.3, 1.2, 5.8, 0.35, fill_rgb=AZUL)
textbox(slide, "💉  LA PROPUESTA", 0.5, 1.22, 5.4, 0.3, size=11, bold=True, color=BLANCO)

componentes = [
    ("Crédito BCI (Opción B)", "$10,000,000", "24 cuotas — $505,611/mes — tasa 1.51%"),
    ("Aporte familiar hermano", "$2,200,000", "Sin interés, sin cuota mensual"),
    ("TOTAL CAPITAL DISPONIBLE", "$12,200,000", "Para pagar las deudas más caras"),
]
for i, (nombre, monto, det) in enumerate(componentes):
    y = 1.65 + i * 0.95
    bg = AZUL_OSCURO if i == 2 else PANEL
    rect(slide, 0.4, y, 5.6, 0.85, fill_rgb=bg)
    textbox(slide, nombre, 0.55, y + 0.05, 3.0, 0.35, size=10, bold=(i==2), color=BLANCO)
    textbox(slide, monto, 3.5, y + 0.02, 2.3, 0.4, size=16, bold=True, color=AZUL, align=PP_ALIGN.RIGHT)
    textbox(slide, det, 0.55, y + 0.5, 5.2, 0.28, size=9, color=GRIS, italic=True)

textbox(slide, "Plan de asignación (de mayor a menor tasa):", 0.5, 4.6, 5.4, 0.3, size=10, bold=True, color=BLANCO)
asignaciones = [
    "✅ CMR Falabella ($1.6M al 3.3%) → libera $80,000/mes",
    "✅ Líneas de crédito ($2.3M al 3.1%) → libera $71,660/mes",
    "✅ Banco Estado ($5.6M al 3.1%) → libera $174,437/mes",
    "⚡ Santander — pago parcial con el remanente",
]
for i, a in enumerate(asignaciones):
    textbox(slide, a, 0.5, 4.95 + i*0.42, 5.5, 0.38, size=10, color=VERDE if "✅" in a else AMBAR)

# Derecha — impacto
rect(slide, 6.5, 1.2, 6.5, 5.8, fill_rgb=PANEL, line_rgb=VERDE, line_width=1)
rect(slide, 6.5, 1.2, 6.5, 0.35, fill_rgb=VERDE)
textbox(slide, "📊  IMPACTO FINANCIERO", 6.7, 1.22, 6.1, 0.3, size=11, bold=True, color=BLANCO)

impactos = [
    ("Cuotas liberadas/mes", "+$326,097", VERDE),
    ("Nueva cuota BCI/mes", "-$505,611", ROJO),
    ("Impacto neto cuotas (corto plazo)", "-$179,514", AMBAR),
    ("Ahorro en intereses/mes", "+$280,268", VERDE),
    ("Ahorro total (24 meses)", "+$3,363,000", VERDE),
    ("Arbitraje de tasa logrado", "1.70%/mes", AZUL),
]
for i, (lbl, val, col) in enumerate(impactos):
    y = 1.65 + i * 0.75
    rect(slide, 6.6, y, 6.3, 0.65, fill_rgb=RGBColor(0x21, 0x26, 0x2D))
    textbox(slide, lbl, 6.75, y + 0.05, 4.5, 0.3, size=10, color=GRIS)
    textbox(slide, val, 10.5, y + 0.03, 2.3, 0.4, size=16, bold=True, color=col, align=PP_ALIGN.RIGHT)

rect(slide, 6.6, 6.2, 6.3, 0.7, fill_rgb=RGBColor(0x1B, 0x2D, 0x1B), line_rgb=VERDE, line_width=1)
textbox(slide, "⚠  Solo ejecutar si se renegocia el alquiler ANTES y se acuerda plan de devolución familiar por escrito.", 6.75, 6.25, 6.1, 0.6, size=10, bold=True, color=VERDE)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — PLAN DE ACCIÓN
# ══════════════════════════════════════════════════════════════════════════════
slide = add_slide()
slide_bg(slide, FONDO)
header_bar(slide, "Plan de Acción Priorizado", "10 acciones en 3 horizontes de tiempo — Impacto acumulado: +$716,000/mes")
footer(slide)

# Horizonte 1
rect(slide, 0.3, 1.25, 12.7, 0.35, fill_rgb=ROJO)
textbox(slide, "HORIZONTE 1 — URGENTE  (0 a 3 meses)  |  Impacto potencial: +$716,000/mes", 0.5, 1.28, 12, 0.28, size=11, bold=True, color=BLANCO)

h1_acciones = [
    ("1", "Renegociar alquiler del taller", "$700,000 → $400,000", "+$300,000/mes", ROJO),
    ("2", "Separar gastos personales de la caja", "~$80K/mes mezclados actualmente", "+$80,000/mes", ROJO),
    ("3", "Renegociar TCs antes de mora formal", "Reducir cuota y tasa de interés", "+$150,000/mes", AMBAR),
    ("4", "Subir precios 10-15% en productos", "Margen bruto pasa de 45% a 50%+", "+$186,000/mes", AMBAR),
]
for i, (num, acc, det, imp, col) in enumerate(h1_acciones):
    x = 0.3 + (i % 2) * 6.4
    y = 1.65 + (i // 2) * 0.72
    rect(slide, x, y, 6.2, 0.65, fill_rgb=PANEL, line_rgb=col, line_width=1)
    textbox(slide, num + ". " + acc, x + 0.15, y + 0.05, 4.0, 0.3, size=11, bold=True, color=BLANCO)
    textbox(slide, det, x + 0.15, y + 0.35, 4.0, 0.25, size=9, color=GRIS, italic=True)
    textbox(slide, imp, x + 4.2, y + 0.1, 1.8, 0.4, size=13, bold=True, color=VERDE, align=PP_ALIGN.RIGHT)

# Horizonte 2
rect(slide, 0.3, 3.18, 12.7, 0.32, fill_rgb=AMBAR)
textbox(slide, "HORIZONTE 2 — IMPORTANTE  (3 a 12 meses)", 0.5, 3.2, 12, 0.28, size=11, bold=True, color=BLANCO)

h2_acciones = [
    ("5", "Escalar ventas a $3,500,000/mes", "Canal MercadoLibre + vendedor comisionista"),
    ("6", "Formalizar empresa (EIRL/SpA)", "Acceso SERCOTEC y FOGAPE — créditos al 0.5-1%/mes"),
    ("7", "Evaluar venta del camión Foton", "Libera $329,778/mes si ventas no suben"),
]
for i, (num, acc, det) in enumerate(h2_acciones):
    x = 0.3 + i * 4.2
    rect(slide, x, 3.55, 4.0, 0.85, fill_rgb=PANEL, line_rgb=AMBAR, line_width=1)
    textbox(slide, num + ". " + acc, x + 0.15, 3.6, 3.7, 0.35, size=10, bold=True, color=BLANCO)
    textbox(slide, det, x + 0.15, 3.95, 3.7, 0.4, size=9, color=GRIS, italic=True)

# Horizonte 3
rect(slide, 0.3, 4.5, 12.7, 0.32, fill_rgb=AZUL)
textbox(slide, "HORIZONTE 3 — LARGO PLAZO  (12 a 36 meses)", 0.5, 4.52, 12, 0.28, size=11, bold=True, color=BLANCO)

h3_acciones = [
    ("8", "Liquidar TCs (CMR primero)", "Libera flujo permanente y mejora historial"),
    ("9", "Colchón liquidez 2 meses (~$3.8M)", "Estabilidad ante imprevistos y temporadas bajas"),
    ("10", "Contratar vendedor/a a comisión", "Multiplica ventas sin aumentar costos fijos"),
]
for i, (num, acc, det) in enumerate(h3_acciones):
    x = 0.3 + i * 4.2
    rect(slide, x, 4.87, 4.0, 0.85, fill_rgb=PANEL, line_rgb=AZUL, line_width=1)
    textbox(slide, num + ". " + acc, x + 0.15, 4.92, 3.7, 0.35, size=10, bold=True, color=BLANCO)
    textbox(slide, det, x + 0.15, 5.27, 3.7, 0.4, size=9, color=GRIS, italic=True)

# Impacto total
rect(slide, 0.3, 5.85, 12.7, 1.0, fill_rgb=RGBColor(0x1B, 0x2D, 0x1B), line_rgb=VERDE, line_width=1)
textbox(slide, "Si se ejecutan las 4 acciones urgentes: +$716,000/mes  |  Déficit actual: -$105,770/mes  →  Superávit potencial: +$610,230/mes", 0.5, 5.95, 12.3, 0.7, size=13, bold=True, color=VERDE, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — DECISIÓN: SEGUIR O CERRAR
# ══════════════════════════════════════════════════════════════════════════════
slide = add_slide()
slide_bg(slide, FONDO)
header_bar(slide, "Decisión Estratégica: ¿Seguir o Cerrar?", "Análisis objetivo de ambos caminos — Basado en datos reales")
footer(slide)

# Columna SEGUIR
rect(slide, 0.3, 1.2, 5.9, 5.8, fill_rgb=PANEL, line_rgb=VERDE, line_width=2)
rect(slide, 0.3, 1.2, 5.9, 0.42, fill_rgb=VERDE)
textbox(slide, "✅  SEGUIR — Condiciones y fundamentos", 0.5, 1.24, 5.5, 0.35, size=12, bold=True, color=BLANCO)

pros = [
    ("Margen bruto sólido", "45-55% por producto — el negocio tiene capacidad real de ganancia"),
    ("Activo productivo real", "Camión Foton + herramientas + know-how de fabricación"),
    ("Potencial de ventas", "Dic-2025 = $3,024,913 (demanda existe, el problema es consistencia)"),
    ("Costo de cierre alto", "Liquidar maquinaria, perder clientes activos, riesgo reputacional"),
    ("Acciones de impacto inmediato", "Renegociar alquiler + separar gastos = +$380K/mes sin vender más"),
    ("Apoyo familiar disponible", "Inyección $2.2M sin costo, red de apoyo activa"),
]
for i, (titulo, desc) in enumerate(pros):
    y = 1.72 + i * 0.82
    rect(slide, 0.4, y, 0.3, 0.3, fill_rgb=VERDE)
    textbox(slide, titulo, 0.8, y, 5.2, 0.3, size=10, bold=True, color=BLANCO)
    textbox(slide, desc, 0.8, y + 0.33, 5.2, 0.38, size=9, color=GRIS, italic=True)

# Columna CERRAR
rect(slide, 6.6, 1.2, 5.9, 5.8, fill_rgb=PANEL, line_rgb=ROJO, line_width=2)
rect(slide, 6.6, 1.2, 5.9, 0.42, fill_rgb=ROJO)
textbox(slide, "⚠  CERRAR — Cuándo considerarlo seriamente", 6.8, 1.24, 5.5, 0.35, size=12, bold=True, color=BLANCO)

cons = [
    ("Ventas no suben en 6 meses", "Si no se supera $2.5M/mes para Sep-2026, la deuda se vuelve impagable"),
    ("Alquiler no se renegocia", "Sin reducir el alquiler, el PE sigue fuera de alcance realista"),
    ("Mora bancaria formal", "Si las TCs entran en cobranza judicial, costos se multiplican"),
    ("Deuda familiar impagable", "Si la hermana no puede devolver, el conflicto familiar es el costo más alto"),
    ("Salud del operador", "El estrés financiero sostenido afecta la calidad del trabajo y las ventas"),
]
for i, (titulo, desc) in enumerate(cons):
    y = 1.72 + i * 0.98
    rect(slide, 6.7, y, 0.3, 0.3, fill_rgb=ROJO)
    textbox(slide, titulo, 7.1, y, 5.2, 0.3, size=10, bold=True, color=BLANCO)
    textbox(slide, desc, 7.1, y + 0.33, 5.2, 0.5, size=9, color=GRIS, italic=True)

# Veredicto
rect(slide, 0.3, 6.28, 12.7, 0.9, fill_rgb=AZUL_OSCURO, line_rgb=AZUL, line_width=1.5)
textbox(slide, "VEREDICTO: El negocio tiene viabilidad CONDICIONADA. Seguir tiene sentido SOLO si en los próximos 90 días se renegocia el alquiler y se separan las finanzas personales. Sin esas dos acciones, el cierre ordenado en 6 meses es la opción más responsable para proteger el patrimonio familiar.", 0.5, 6.33, 12.3, 0.8, size=11, bold=True, color=BLANCO, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — CIERRE Y PRÓXIMOS PASOS
# ══════════════════════════════════════════════════════════════════════════════
slide = add_slide()
slide_bg(slide, FONDO)
header_bar(slide, "Próximos Pasos — Los 90 días críticos", "Las 3 acciones que definen si el negocio sobrevive")
footer(slide)

pasos = [
    ("ESTA SEMANA", "Conversación con el arrendador del taller",
     "Presentar propuesta: $400,000/mes con contrato firmado. Si dice que no → buscar alternativa de local.",
     ROJO, "1"),
    ("ESTE MES", "Separar caja personal de caja del negocio",
     "Abrir cuenta bancaria exclusiva para el negocio. Todo ingreso del negocio entra ahí. Todo gasto personal sale de otra cuenta.",
     AMBAR, "2"),
    ("PRÓXIMOS 30 DÍAS", "Negociar TCs con los bancos",
     "Llamar a Santander y CMR para renegociar antes de entrar en mora. Objetivo: reducir cuota total en $150,000/mes.",
     AZUL, "3"),
]

for i, (plazo, accion, detalle, col, num) in enumerate(pasos):
    y = 1.25 + i * 1.85
    rect(slide, 0.3, y, 12.7, 1.7, fill_rgb=PANEL, line_rgb=col, line_width=1.5)
    rect(slide, 0.3, y, 1.2, 1.7, fill_rgb=col)
    textbox(slide, num, 0.55, y + 0.55, 0.7, 0.7, size=32, bold=True, color=BLANCO, align=PP_ALIGN.CENTER)
    textbox(slide, plazo, 1.65, y + 0.08, 3.5, 0.35, size=10, bold=True, color=col)
    textbox(slide, accion, 1.65, y + 0.38, 10.8, 0.45, size=14, bold=True, color=BLANCO)
    textbox(slide, detalle, 1.65, y + 0.88, 10.8, 0.6, size=10, color=GRIS, italic=True)

rect(slide, 0.3, 6.82, 12.7, 0.45, fill_rgb=PANEL)
textbox(slide, "Meta a 90 días: Déficit actual -$105,770/mes  →  Superávit estimado +$374,230/mes  |  PE bajará de $4,250,896 a $3,582,484", 0.5, 6.87, 12.3, 0.35, size=11, bold=True, color=VERDE, align=PP_ALIGN.CENTER)

output = r"C:\ClaudeWork\chiquito_financiero\Presentacion_ChiquitoFinanzas.pptx"
prs.save(output)
print(f"[OK] PowerPoint guardado: {output}")
