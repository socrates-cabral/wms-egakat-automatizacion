import sys
try:
    sys.stdout.reconfigure(encoding="utf-8")
except Exception:
    pass

# pptx_report.py — Generador de presentación ejecutiva PowerPoint
# Llamado desde main.py con los datos reales de la sesión

import io
import pandas as pd
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── Paleta ─────────────────────────────────────────────────────────────────────
AZUL        = RGBColor(0x1F, 0x6F, 0xEB)
AZUL_OSCURO = RGBColor(0x0D, 0x47, 0xA1)
VERDE       = RGBColor(0x1E, 0x8B, 0x4C)
ROJO        = RGBColor(0xC0, 0x39, 0x2B)
AMBAR       = RGBColor(0x9A, 0x7D, 0x0A)
GRIS        = RGBColor(0x5D, 0x6D, 0x7E)
GRIS_CLARO  = RGBColor(0xF2, 0xF3, 0xF4)
BLANCO      = RGBColor(0xFF, 0xFF, 0xFF)
FONDO       = RGBColor(0x0D, 0x11, 0x17)
PANEL       = RGBColor(0x16, 0x1B, 0x22)


def _fmt(valor) -> str:
    """Formato chileno: $1.234.567"""
    try:
        return "$" + f"{float(valor):,.0f}".replace(",", ".")
    except (TypeError, ValueError):
        return "—"


def _fmt_no_signo(valor) -> str:
    try:
        return f"{float(valor):,.0f}".replace(",", ".")
    except (TypeError, ValueError):
        return "—"


# ── Helpers de dibujo ──────────────────────────────────────────────────────────
def _blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def _bg(slide, color=FONDO):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def _rect(slide, l, t, w, h, fill=None, line=None, lw=None):
    s = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    s.line.fill.background()
    if fill:
        s.fill.solid(); s.fill.fore_color.rgb = fill
    else:
        s.fill.background()
    if line:
        s.line.color.rgb = line
        if lw:
            s.line.width = Pt(lw)
    else:
        s.line.fill.background()
    return s


def _txt(slide, text, l, t, w, h, size=12, bold=False, color=BLANCO,
         align=PP_ALIGN.LEFT, italic=False):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = color


def _header(slide, titulo, subtitulo=""):
    _rect(slide, 0, 0, 13.33, 1.05, fill=PANEL)
    _rect(slide, 0, 1.0, 13.33, 0.05, fill=AZUL)
    _txt(slide, titulo, 0.3, 0.1, 9.5, 0.55, size=22, bold=True, color=BLANCO)
    if subtitulo:
        _txt(slide, subtitulo, 0.3, 0.62, 9.5, 0.33, size=11, color=GRIS, italic=True)
    _txt(slide, "Chiquito Finanzas | Mar 2026", 10.3, 0.1, 2.9, 0.35, size=9,
         color=GRIS, align=PP_ALIGN.RIGHT)


def _footer(slide):
    _rect(slide, 0, 7.2, 13.33, 0.3, fill=PANEL)
    _txt(slide, "Chiquito Finanzas  |  Diagnóstico Financiero  |  Confidencial",
         0.3, 7.22, 12, 0.25, size=8, color=GRIS)


def _kpi(slide, l, t, valor, label, color=AZUL, w=2.5, h=1.1):
    _rect(slide, l, t, w, h, fill=PANEL, line=color, lw=1.5)
    _txt(slide, valor, l+0.1, t+0.07, w-0.2, 0.52, size=20, bold=True,
         color=color, align=PP_ALIGN.CENTER)
    _txt(slide, label, l+0.1, t+0.62, w-0.2, 0.36, size=9,
         color=GRIS, align=PP_ALIGN.CENTER)


def _tabla(slide, headers, rows, l, t, w, h):
    nc, nr = len(headers), len(rows) + 1
    tbl = slide.shapes.add_table(nr, nc, Inches(l), Inches(t),
                                  Inches(w), Inches(h)).table
    cw = Inches(w) // nc
    for i in range(nc):
        tbl.columns[i].width = cw
    for ci, hdr in enumerate(headers):
        cell = tbl.cell(0, ci)
        cell.text = hdr
        cell.fill.solid(); cell.fill.fore_color.rgb = AZUL
        for p in cell.text_frame.paragraphs:
            for r in p.runs:
                r.font.bold = True; r.font.color.rgb = BLANCO; r.font.size = Pt(10)
            p.alignment = PP_ALIGN.CENTER
    for ri, row in enumerate(rows):
        for ci, val in enumerate(row):
            cell = tbl.cell(ri+1, ci)
            cell.text = str(val)
            cell.fill.solid()
            cell.fill.fore_color.rgb = GRIS_CLARO if ri % 2 == 0 else BLANCO
            for p in cell.text_frame.paragraphs:
                for r in p.runs:
                    r.font.size = Pt(9)


# ══════════════════════════════════════════════════════════════════════════════
# FUNCIÓN PRINCIPAL
# ══════════════════════════════════════════════════════════════════════════════
def generar_presentacion(kpis: dict, df_resumen: pd.DataFrame,
                         df_deuda: pd.DataFrame) -> bytes:
    """
    Genera la presentación ejecutiva PowerPoint con los datos reales de la sesión.

    kpis: dict con claves prom_ing, prom_gas, prom_res, total_deuda,
                          total_cuotas, pe_actual, pct_pe, pct_cuotas,
                          instrumentos_activos
    df_resumen: columnas [mes, ingresos, gastos, resultado]
    df_deuda:   columnas [acreedor, saldo, cuota, tasa, tipo]
    Retorna bytes del archivo .pptx
    """
    prs = Presentation()
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.5)

    pi   = kpis.get('prom_ing', 0)
    pg   = kpis.get('prom_gas', 0)
    pr   = kpis.get('prom_res', 0)
    td   = kpis.get('total_deuda', 0)
    tc   = kpis.get('total_cuotas', 0)
    pe   = kpis.get('pe_actual', 0)
    pct  = kpis.get('pct_pe', 0)
    pctc = kpis.get('pct_cuotas', 0)
    inst = kpis.get('instrumentos_activos', 0)

    color_res = ROJO if pr < 0 else VERDE
    color_pe  = ROJO if pct < 50 else (AMBAR if pct < 70 else VERDE)

    # ── SLIDE 1: Portada ──────────────────────────────────────────────────────
    slide = _blank(prs)
    _bg(slide)
    _rect(slide, 0, 0, 13.33, 0.08, fill=AZUL)
    _rect(slide, 1.5, 1.8, 10.3, 4.0, fill=PANEL, line=AZUL, lw=1)
    _rect(slide, 1.5, 1.8, 10.3, 0.08, fill=AZUL)
    _txt(slide, "DIAGNÓSTICO FINANCIERO", 2.0, 2.15, 9.5, 0.5,
         size=12, bold=True, color=GRIS, align=PP_ALIGN.CENTER)
    _txt(slide, "Chiquito Finanzas", 2.0, 2.55, 9.5, 0.9,
         size=38, bold=True, color=BLANCO, align=PP_ALIGN.CENTER)
    _txt(slide, "Taller de Muebles  |  Macul, Santiago, Chile",
         2.0, 3.45, 9.5, 0.45, size=15, color=GRIS, align=PP_ALIGN.CENTER, italic=True)
    _rect(slide, 4.5, 4.1, 4.3, 0.04, fill=AZUL)
    estado_txt = "Estado: DÉFICIT OPERATIVO — Requiere acción inmediata" if pr < 0 \
                 else "Estado: OPERANDO — Monitoreo activo recomendado"
    _txt(slide, estado_txt, 2.0, 4.25, 9.5, 0.4,
         size=12, bold=True, color=ROJO if pr < 0 else VERDE, align=PP_ALIGN.CENTER)
    _txt(slide, "Presentado por: Sócrates Cabral  |  Control de Gestión y Mejora Continua  |  Egakat SPA",
         2.0, 4.75, 9.5, 0.35, size=10, color=GRIS, align=PP_ALIGN.CENTER)
    _rect(slide, 0, 7.42, 13.33, 0.08, fill=AZUL)

    # ── SLIDE 2: Resumen ejecutivo ────────────────────────────────────────────
    slide = _blank(prs)
    _bg(slide)
    _header(slide, "Resumen Ejecutivo", "El diagnóstico en 60 segundos")
    _footer(slide)

    _rect(slide, 0.3, 1.2, 5.8, 5.8, fill=PANEL, line=ROJO, lw=1)
    _rect(slide, 0.3, 1.2, 5.8, 0.35, fill=ROJO)
    _txt(slide, "SITUACIÓN ACTUAL", 0.5, 1.24, 5.4, 0.28, size=10, bold=True)

    datos_sit = [
        f"Ingresos promedio: {_fmt(pi)}/mes",
        f"Gastos promedio: {_fmt(pg)}/mes",
        f"Resultado: {_fmt(pr)}/mes ({'déficit' if pr < 0 else 'superávit'})",
        f"Deuda total: {_fmt(td)}",
        f"Cuotas mensuales: {_fmt(tc)} ({pctc:.0f}% del ingreso)",
        f"Punto de Equilibrio: {_fmt(pe)}/mes",
        f"Ventas actuales = {pct:.0f}% del PE",
    ]
    for i, d in enumerate(datos_sit):
        col = BLANCO if i < 3 else GRIS
        _txt(slide, "• " + d, 0.5, 1.65 + i*0.65, 5.4, 0.55, size=10, color=col)

    _rect(slide, 6.5, 1.2, 6.5, 5.8, fill=PANEL, line=AMBAR, lw=1)
    _rect(slide, 6.5, 1.2, 6.5, 0.35, fill=AMBAR)
    _txt(slide, "5 CAUSAS DEL DÉFICIT", 6.7, 1.24, 6.1, 0.28, size=10, bold=True)

    causas = [
        ("1.", "Deuda aplastante", f"Cuotas = {pctc:.0f}% del ingreso promedio"),
        ("2.", "Alquiler excesivo", "$700K/mes = 38% del ingreso (normal: 10-15%)"),
        ("3.", "Ventas insuficientes", f"Vende {_fmt(pi)} vs. necesita {_fmt(pe)}"),
        ("4.", "Gastos personales mezclados", "~$80K/mes de gastos privados en la caja"),
        ("5.", "TCs en mora o riesgo de mora", "Genera tasa TMC 2.75%/mes — la más alta"),
    ]
    for i, (num, tit, desc) in enumerate(causas):
        y = 1.65 + i * 1.0
        _txt(slide, num, 6.7, y, 0.4, 0.35, size=18, bold=True, color=AMBAR)
        _txt(slide, tit, 7.1, y, 5.5, 0.35, size=12, bold=True)
        _txt(slide, desc, 7.1, y+0.38, 5.5, 0.4, size=10, color=GRIS, italic=True)

    # ── SLIDE 3: KPIs ─────────────────────────────────────────────────────────
    slide = _blank(prs)
    _bg(slide)
    _header(slide, "KPIs Financieros",
            f"Datos reales del libro de caja — {len(df_resumen)} meses registrados")
    _footer(slide)

    kpi_row1 = [
        (_fmt(pi),  "Ingreso prom/mes",    VERDE),
        (_fmt(pg),  "Gasto prom/mes",      ROJO),
        (_fmt(pr),  "Resultado neto prom", color_res),
        (f"{pct:.0f}%", "% PE alcanzado", color_pe),
    ]
    for i, (v, l, c) in enumerate(kpi_row1):
        _kpi(slide, 0.3 + i*3.15, 1.2, v, l, c, w=3.0, h=1.1)

    kpi_row2 = [
        (_fmt(td), "Deuda total",         ROJO),
        (_fmt(tc), "Cuotas/mes",          AMBAR),
        (f"{pctc:.0f}%", "% ingreso en cuotas", AMBAR),
        (str(inst), "Instrumentos activos", AZUL),
    ]
    for i, (v, l, c) in enumerate(kpi_row2):
        _kpi(slide, 0.3 + i*3.15, 2.5, v, l, c, w=3.0, h=1.1)

    # Tabla flujo mensual con datos reales
    _rect(slide, 0.3, 3.75, 12.7, 3.45, fill=PANEL)
    _txt(slide, "Flujo mensual real", 0.5, 3.82, 8, 0.35, size=12, bold=True)

    filas_tabla = []
    for _, row in df_resumen.iterrows():
        res = float(row.get('resultado', 0))
        estado = "✅ Superávit" if res > 0 else "❌ Déficit"
        filas_tabla.append([
            str(row.get('mes', '')),
            _fmt(row.get('ingresos', 0)),
            _fmt(row.get('gastos', 0)),
            _fmt(res),
            estado,
        ])

    _tabla(slide, ["Mes", "Ingresos", "Gastos", "Resultado", "Estado"],
           filas_tabla, 0.3, 4.2, 12.7, min(2.9, 0.45 * len(filas_tabla) + 0.5))

    # ── SLIDE 4: Deuda ────────────────────────────────────────────────────────
    slide = _blank(prs)
    _bg(slide)
    _header(slide, "Estructura de Deuda",
            f"Total: {_fmt(td)}  |  Cuotas: {_fmt(tc)}/mes  |  {inst} instrumentos activos")
    _footer(slide)

    filas_deuda = []
    if not df_deuda.empty:
        for _, row in df_deuda.iterrows():
            saldo = float(pd.to_numeric(row.get('saldo', 0), errors='coerce') or 0)
            cuota = float(pd.to_numeric(row.get('cuota', 0), errors='coerce') or 0)
            tasa  = float(pd.to_numeric(row.get('tasa',  0), errors='coerce') or 0)
            if saldo > 0 or cuota > 0:
                filas_deuda.append([
                    str(row.get('acreedor', '')),
                    _fmt(saldo),
                    _fmt(cuota),
                    f"{tasa:.1f}%/mes" if tasa > 0 else "—",
                    str(row.get('tipo', '')).title(),
                ])

    if filas_deuda:
        _tabla(slide,
               ["Acreedor", "Saldo", "Cuota/mes", "Tasa", "Tipo"],
               filas_deuda, 0.3, 1.2, 12.7,
               min(5.5, 0.5 * len(filas_deuda) + 0.6))

    _rect(slide, 0.3, 6.65, 12.7, 0.75, fill=RGBColor(0x2D, 0x1B, 0x1B),
          line=ROJO, lw=1)
    _txt(slide,
         f"El {pctc:.0f}% del ingreso mensual se destina a pagar deudas. "
         "Sin reducir la deuda o aumentar ventas significativamente, el negocio no es viable.",
         0.5, 6.7, 12.3, 0.65, size=11, bold=True, color=ROJO)

    # ── SLIDE 5: Punto de equilibrio ──────────────────────────────────────────
    slide = _blank(prs)
    _bg(slide)
    _header(slide, "Punto de Equilibrio",
            f"Ventas actuales {_fmt(pi)}/mes = {pct:.0f}% del PE ({_fmt(pe)}/mes)")
    _footer(slide)

    _rect(slide, 0.3, 1.2, 12.7, 2.3, fill=PANEL)
    _txt(slide, f"VENTAS ACTUALES: {_fmt(pi)}/mes", 0.5, 1.28, 6, 0.38,
         size=13, bold=True, color=VERDE)
    _txt(slide, f"PUNTO DE EQUILIBRIO: {_fmt(pe)}/mes", 6.5, 1.28, 6.3, 0.38,
         size=13, bold=True, color=AMBAR)

    # Barra progreso
    _rect(slide, 0.5, 1.8, 12.3, 0.55, fill=RGBColor(0x21, 0x26, 0x2D))
    barra_w = min(12.3 * pct / 100, 12.3)
    col_barra = ROJO if pct < 50 else (AMBAR if pct < 70 else VERDE)
    _rect(slide, 0.5, 1.8, max(barra_w, 0.3), 0.55, fill=col_barra)
    _txt(slide, f"{pct:.0f}%", 0.5 + barra_w/2 - 0.3, 1.83, 0.7, 0.48,
         size=16, bold=True, align=PP_ALIGN.CENTER)
    _txt(slide, "Ventas actuales", 0.5, 2.4, 5, 0.28, size=9, color=GRIS)
    _txt(slide, "Punto de Equilibrio (100%)", 8.5, 2.4, 4, 0.28,
         size=9, color=GRIS, align=PP_ALIGN.RIGHT)

    escenarios = [
        ["Actual (sin cambios)", _fmt(pe), "Solo aumentar ventas", f"{pct:.0f}% → 100%", "Muy difícil"],
        ["Renegociar alquiler", "$3.582.484", "Alquiler $700K → $400K", f"{pct:.0f}% → 100%", "Difícil pero posible"],
        ["Renegociar + subir precio", "$3.186.484", "Alquiler $400K + precios +15%", f"{pct:.0f}% → 100%", "Alcanzable 6-12 m"],
        ["Inyección + renegociar", "$2.962.484", "Crédito BCI + alquiler $400K", f"{pct:.0f}% → 100%", "Alcanzable 3-6 m"],
    ]
    _tabla(slide, ["Escenario", "Ventas necesarias", "Cambios", "% PE", "Factibilidad"],
           escenarios, 0.3, 3.75, 12.7, 3.1)

    # ── SLIDE 6: Plan de acción ────────────────────────────────────────────────
    slide = _blank(prs)
    _bg(slide)
    _header(slide, "Plan de Acción Priorizado",
            "10 acciones en 3 horizontes — Impacto acumulado potencial: +$716.000/mes")
    _footer(slide)

    _rect(slide, 0.3, 1.2, 12.7, 0.32, fill=ROJO)
    _txt(slide, "HORIZONTE 1 — URGENTE (0-3 meses)  |  +$716.000/mes potencial",
         0.5, 1.23, 12, 0.27, size=10, bold=True)

    h1 = [
        ("1", "Renegociar alquiler taller $700K → $400K", "+$300.000/mes", ROJO),
        ("2", "Separar gastos personales de la caja", "+$80.000/mes", ROJO),
        ("3", "Renegociar TCs antes de mora formal", "+$150.000/mes", AMBAR),
        ("4", "Subir precios 10-15%", "+$186.000/mes", AMBAR),
    ]
    for i, (n, acc, imp, col) in enumerate(h1):
        x = 0.3 + (i % 2) * 6.4
        y = 1.57 + (i // 2) * 0.72
        _rect(slide, x, y, 6.2, 0.65, fill=PANEL, line=col, lw=1)
        _txt(slide, f"{n}. {acc}", x+0.15, y+0.05, 4.2, 0.3, size=10, bold=True)
        _txt(slide, imp, x+4.3, y+0.1, 1.7, 0.38, size=12, bold=True,
             color=VERDE, align=PP_ALIGN.RIGHT)

    _rect(slide, 0.3, 3.1, 12.7, 0.3, fill=AMBAR)
    _txt(slide, "HORIZONTE 2 — IMPORTANTE (3-12 meses)", 0.5, 3.12, 12, 0.26,
         size=10, bold=True)
    h2 = [
        ("5", "Escalar ventas a $3.500.000/mes", "Meta supervivencia"),
        ("6", "Formalizar empresa (EIRL/SpA)", "Acceso SERCOTEC/FOGAPE"),
        ("7", "Evaluar venta camión Foton", "Libera $329.778/mes"),
    ]
    for i, (n, acc, imp) in enumerate(h2):
        x = 0.3 + i * 4.2
        _rect(slide, x, 3.45, 4.0, 0.82, fill=PANEL, line=AMBAR, lw=1)
        _txt(slide, f"{n}. {acc}", x+0.15, 3.5, 3.7, 0.35, size=10, bold=True)
        _txt(slide, imp, x+0.15, 3.85, 3.7, 0.35, size=9, color=GRIS, italic=True)

    _rect(slide, 0.3, 4.37, 12.7, 0.3, fill=AZUL)
    _txt(slide, "HORIZONTE 3 — LARGO PLAZO (12-36 meses)", 0.5, 4.39, 12, 0.26,
         size=10, bold=True)
    h3 = [
        ("8", "Liquidar TCs (CMR primero)", "Libera flujo permanente"),
        ("9", "Colchón liquidez 2 meses (~$3.800.000)", "Estabilidad"),
        ("10", "Vendedor/a a comisión", "Multiplica ventas sin costo fijo"),
    ]
    for i, (n, acc, imp) in enumerate(h3):
        x = 0.3 + i * 4.2
        _rect(slide, x, 4.72, 4.0, 0.82, fill=PANEL, line=AZUL, lw=1)
        _txt(slide, f"{n}. {acc}", x+0.15, 4.77, 3.7, 0.35, size=10, bold=True)
        _txt(slide, imp, x+0.15, 5.12, 3.7, 0.35, size=9, color=GRIS, italic=True)

    _rect(slide, 0.3, 5.68, 12.7, 0.88, fill=RGBColor(0x1B, 0x2D, 0x1B),
          line=VERDE, lw=1)
    _txt(slide,
         f"Si se ejecutan las 4 acciones urgentes: resultado estimado +$610.230/mes  "
         f"(déficit actual {_fmt(pr)}/mes → superávit potencial +$610.230/mes)",
         0.5, 5.75, 12.3, 0.75, size=11, bold=True, color=VERDE,
         align=PP_ALIGN.CENTER)

    # ── SLIDE 7: Decisión seguir o cerrar ─────────────────────────────────────
    slide = _blank(prs)
    _bg(slide)
    _header(slide, "Decisión: ¿Seguir o Cerrar?",
            "Análisis objetivo de ambos caminos basado en datos reales")
    _footer(slide)

    _rect(slide, 0.3, 1.2, 5.9, 5.8, fill=PANEL, line=VERDE, lw=2)
    _rect(slide, 0.3, 1.2, 5.9, 0.4, fill=VERDE)
    _txt(slide, "SEGUIR — Fundamentos", 0.5, 1.24, 5.5, 0.32, size=11, bold=True)
    pros = [
        ("Margen bruto sólido", "45-55% — el negocio tiene capacidad real de ganancia"),
        ("Activo productivo real", "Camión + herramientas + know-how acumulado"),
        ("Demanda demostrada", f"Dic-2025 = $3.024.913 — la demanda existe"),
        ("Acciones de impacto inmediato", "Renegociar alquiler = +$300K/mes sin vender más"),
        ("Apoyo familiar disponible", "Inyección $2.2M sin costo, red activa"),
    ]
    for i, (tit, desc) in enumerate(pros):
        y = 1.7 + i * 0.92
        _rect(slide, 0.4, y, 0.28, 0.28, fill=VERDE)
        _txt(slide, tit, 0.8, y, 5.2, 0.3, size=10, bold=True)
        _txt(slide, desc, 0.8, y+0.32, 5.2, 0.45, size=9, color=GRIS, italic=True)

    _rect(slide, 6.6, 1.2, 5.9, 5.8, fill=PANEL, line=ROJO, lw=2)
    _rect(slide, 6.6, 1.2, 5.9, 0.4, fill=ROJO)
    _txt(slide, "CERRAR — Cuándo considerarlo", 6.8, 1.24, 5.5, 0.32, size=11, bold=True)
    contras = [
        ("Ventas no suben en 6 meses", "Si no supera $2.5M/mes para Sep-2026"),
        ("Alquiler no se renegocia", "Sin reducirlo, el PE sigue fuera de alcance"),
        ("Mora bancaria formal", "Si TCs entran en cobranza, costos se multiplican"),
        ("Deuda familiar impagable", "El conflicto familiar es el costo más alto"),
        ("Desgaste del operador", "El estrés sostenido afecta la calidad y ventas"),
    ]
    for i, (tit, desc) in enumerate(contras):
        y = 1.7 + i * 0.92
        _rect(slide, 6.7, y, 0.28, 0.28, fill=ROJO)
        _txt(slide, tit, 7.1, y, 5.1, 0.3, size=10, bold=True)
        _txt(slide, desc, 7.1, y+0.32, 5.1, 0.45, size=9, color=GRIS, italic=True)

    _rect(slide, 0.3, 6.25, 12.7, 0.88, fill=AZUL_OSCURO, line=AZUL, lw=1.5)
    veredicto = (
        "VEREDICTO: Viabilidad CONDICIONADA. Seguir tiene sentido SOLO si en 90 días "
        "se renegocia el alquiler y se separan finanzas personales. "
        "Sin esas dos acciones, el cierre ordenado en 6 meses es la opción más responsable."
    )
    _txt(slide, veredicto, 0.5, 6.32, 12.3, 0.78, size=11, bold=True,
         align=PP_ALIGN.CENTER)

    # ── Serializar a bytes ────────────────────────────────────────────────────
    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf.read()
